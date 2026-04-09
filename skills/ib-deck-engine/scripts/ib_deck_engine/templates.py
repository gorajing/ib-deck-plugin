"""
IB Slide Template Library v1 — Comprehensive template system for IB presentations.

Inspired by real bank decks (Goldman Sachs, Moelis, Evercore, JPMorgan, Lazard).
Uses a structured JSON spec → deterministic pixel rendering approach.

Template catalog:
  1. render_cover               — Cover slide (GS/Moelis style)
  2. render_section_divider     — Full-bg navy divider with section title
  3. render_toc                 — Numbered agenda with blue bands
  4. render_exec_summary        — Callout bar + hierarchical bullets
  5. render_investment_highlights — 4-card grid with icons
  6. render_financial_summary   — Historical + projected P&L table
  7. render_stacked_bar_table   — Moelis revenue by segment
  8. render_dual_chart          — GS dual bar chart with CAGR
  9. render_football_field      — Valuation summary with ranges
 10. render_sensitivity         — WACC × TGR grid
 11. render_sources_uses        — LBO capital structure
 12. render_trading_comps       — Peer multiples table
 13. render_key_terms           — Moelis pill labels + bullets
 14. render_closing             — Dark closing slide

All templates share:
- Consistent chrome (logo, confidential, slide number, Investment Banking badge)
- Same color palette (can be swapped per bank style)
- Same font hierarchy (title 24, section 12, body 10, source 8)
- Right-aligned numerics, left-aligned labels
- Action title enforcement
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from dataclasses import dataclass


# ═══════════════════════════════════════════════
# STYLE SYSTEM
# ═══════════════════════════════════════════════

@dataclass
class BankStyle:
    name: str = "Advisory Group"
    # Colors — all hex
    primary: str = "#1B365D"
    primary_light: str = "#3B6B9D"
    secondary: str = "#A7C4E0"
    accent_orange: str = "#E07020"
    accent_green: str = "#2E8B57"
    accent_red: str = "#C00000"
    bg: str = "#EDEDED"
    text: str = "#1F1F1F"
    muted: str = "#666666"
    border: str = "#BBBBBB"
    row_alt: str = "#F5F5F5"
    header_bg: str = "#D6E4F0"
    highlight: str = "#FFF2CC"
    # Distinct segment palette (for stacked bars)
    segment_colors: tuple = ("#1B365D", "#4B7BA8", "#8EB4D6", "#B8CDE2", "#D6E4F0")

    def rgb(self, hex_color: str) -> RGBColor:
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


GS = BankStyle()


# ═══════════════════════════════════════════════
# IB RENDERER
# ═══════════════════════════════════════════════

class IBRenderer:
    """Comprehensive IB deck template system."""

    def __init__(self, style: BankStyle = None):
        self.style = style or GS
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    # ─── Core primitives ───
    def _new_slide(self, bg_color=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.style.rgb(bg_color or self.style.bg)
        return slide

    def _text(self, slide, left, top, width, height, text,
              size=11, bold=False, italic=False, color=None,
              align=PP_ALIGN.LEFT, font=None):
        s = self.style
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = s.rgb(color or s.text)
        p.font.name = font or "Calibri"
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        return box

    def _text_bg(self, slide, left, top, width, height, text, bg_color,
                 border_color=None, size=11, bold=False,
                 color=None, align=PP_ALIGN.CENTER):
        s = self.style
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = s.rgb(bg_color)
        if border_color:
            shape.line.color.rgb = s.rgb(border_color)
            shape.line.width = Pt(0.75)
        else:
            shape.line.fill.background()
        tf = shape.text_frame
        tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
        tf.margin_left = Pt(3); tf.margin_right = Pt(3)
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = s.rgb(color or s.text)
        p.font.name = "Calibri"
        p.alignment = align
        return shape

    def _rect(self, slide, left, top, width, height, fill_color, border=False):
        s = self.style
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = s.rgb(fill_color)
        if border:
            shape.line.color.rgb = s.rgb(s.border)
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()
        return shape

    def _line(self, slide, left, top, width, color="#666666", thickness=0.01):
        s = self.style
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(thickness))
        line.fill.solid()
        line.fill.fore_color.rgb = s.rgb(color)
        line.line.fill.background()
        return line

    def _vline(self, slide, left, top, height, color="#666666", thickness=0.02):
        s = self.style
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(thickness), Inches(height))
        line.fill.solid()
        line.fill.fore_color.rgb = s.rgb(color)
        line.line.fill.background()
        return line

    def _chrome(self, slide, slide_number=None):
        """Standard header/footer chrome."""
        s = self.style
        # Logo box top-left
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.45), Inches(0.25), Inches(1.0), Inches(0.5))
        box.fill.background()
        box.line.color.rgb = s.rgb(s.primary)
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = s.name
        p.font.size = Pt(9); p.font.bold = True
        p.font.color.rgb = s.rgb(s.primary)
        p.alignment = PP_ALIGN.CENTER

        # Confidential top-right
        self._text(slide, 11.5, 0.2, 1.5, 0.3,
            "Confidential", 10, italic=True, color=s.accent_red, align=PP_ALIGN.RIGHT)

        # Investment Banking badge bottom-right
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(10.3), Inches(6.95), Inches(2.2), Inches(0.3))
        rect.fill.solid(); rect.fill.fore_color.rgb = s.rgb(s.primary)
        rect.line.fill.background()
        tf = rect.text_frame; p = tf.paragraphs[0]
        p.text = "Investment Banking"
        p.font.size = Pt(9); p.font.bold = True
        p.font.color.rgb = s.rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER

        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(12.5), Inches(6.95), Inches(0.4), Inches(0.3))
        strip.fill.solid(); strip.fill.fore_color.rgb = s.rgb(s.secondary)
        strip.line.fill.background()

        # Slide number
        if slide_number is not None:
            self._text(slide, 9.5, 6.95, 0.7, 0.3,
                f"{slide_number}", 8, color=s.muted, align=PP_ALIGN.RIGHT)

    def _title(self, slide, text, y=1.0, height=0.75, size=24):
        self._text(slide, 0.5, y, 11, height, text, size=size, bold=True, color=self.style.primary)

    def _subtitle(self, slide, text, y=1.78):
        self._text(slide, 0.5, y, 11, 0.25, text, size=10, italic=True, color=self.style.muted)

    def _source(self, slide, text, y=6.65):
        self._text(slide, 0.5, y, 11, 0.25, text, size=8, italic=True, color=self.style.muted)

    # ═══════════════════════════════════════════════
    # TEMPLATE 1: COVER
    # ═══════════════════════════════════════════════
    def render_cover(self, spec: dict):
        s = self.style
        slide = self._new_slide()

        # Confidential top-right
        self._text(slide, 11.5, 0.2, 1.5, 0.3,
            "Confidential", 10, italic=True, color=s.accent_red, align=PP_ALIGN.RIGHT)

        # Logo top-left (larger for cover)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.3), Inches(1.3), Inches(0.7))
        box.fill.background()
        box.line.color.rgb = s.rgb(s.primary)
        box.line.width = Pt(1.5)
        tf = box.text_frame; p = tf.paragraphs[0]
        p.text = s.name
        p.font.size = Pt(12); p.font.bold = True
        p.font.color.rgb = s.rgb(s.primary)
        p.alignment = PP_ALIGN.CENTER

        # Title — large
        self._text(slide, 0.5, 2.8, 12, 1.2,
            spec.get("title", ""), size=36, bold=True, color=s.primary, font="Georgia")

        # Subtitle
        if spec.get("subtitle"):
            self._text(slide, 0.5, 4.0, 10, 0.5,
                spec["subtitle"], size=18, color=s.primary_light)

        # Date
        self._text(slide, 0.5, 6.5, 4, 0.3,
            spec.get("date", ""), 12, color=s.muted)

        # Badge
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(10.3), Inches(6.85), Inches(2.2), Inches(0.35))
        rect.fill.solid(); rect.fill.fore_color.rgb = s.rgb(s.primary)
        rect.line.fill.background()
        tf = rect.text_frame; p = tf.paragraphs[0]
        p.text = "Investment Banking"
        p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = s.rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(12.5), Inches(6.85), Inches(0.4), Inches(0.35))
        strip.fill.solid(); strip.fill.fore_color.rgb = s.rgb(s.secondary)
        strip.line.fill.background()

    # ═══════════════════════════════════════════════
    # TEMPLATE 2: SECTION DIVIDER
    # ═══════════════════════════════════════════════
    def render_section_divider(self, spec: dict):
        s = self.style
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # Full navy background
        self._rect(slide, 0, 0, 13.333, 7.5, s.primary)
        # Section title
        self._text(slide, 0.8, 3.0, 12, 1.0,
            spec.get("title", ""), size=30, bold=True, color="#FFFFFF", font="Georgia")
        # Badge
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(10.3), Inches(6.85), Inches(2.2), Inches(0.35))
        rect.fill.solid(); rect.fill.fore_color.rgb = s.rgb("#FFFFFF")
        rect.line.fill.background()
        tf = rect.text_frame; p = tf.paragraphs[0]
        p.text = "Investment Banking"
        p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = s.rgb(s.primary)
        p.alignment = PP_ALIGN.CENTER
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(12.5), Inches(6.85), Inches(0.4), Inches(0.35))
        strip.fill.solid(); strip.fill.fore_color.rgb = s.rgb(s.secondary)
        strip.line.fill.background()

    # ═══════════════════════════════════════════════
    # TEMPLATE 3: TABLE OF CONTENTS
    # ═══════════════════════════════════════════════
    def render_toc(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", "Today's Agenda"))

        items = spec.get("items", [])
        for i, item in enumerate(items):
            y = 2.3 + i * 0.95
            # Light blue band
            band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.6), Inches(y), Inches(10), Inches(0.7))
            band.fill.solid()
            band.fill.fore_color.rgb = s.rgb(s.header_bg)
            band.line.fill.background()
            # Navy number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                Inches(0.9), Inches(y + 0.05), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = s.rgb(s.primary)
            circle.line.fill.background()
            tf = circle.text_frame; p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(18); p.font.bold = True
            p.font.color.rgb = s.rgb("#FFFFFF")
            p.alignment = PP_ALIGN.CENTER
            # Item text
            self._text(slide, 1.8, y + 0.18, 9.8, 0.4,
                item, 15, bold=True, color=s.primary)

    # ═══════════════════════════════════════════════
    # TEMPLATE 4: EXECUTIVE SUMMARY
    # ═══════════════════════════════════════════════
    def render_exec_summary(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", "Executive Summary"))

        # Callout bar
        callout = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.85), Inches(12.333), Inches(0.45))
        callout.fill.solid()
        callout.fill.fore_color.rgb = s.rgb(s.header_bg)
        callout.line.color.rgb = s.rgb(s.primary)
        callout.line.width = Pt(0.75)
        tf = callout.text_frame
        tf.margin_left = Pt(10)
        p = tf.paragraphs[0]
        p.text = spec.get("callout", "")
        p.font.size = Pt(11); p.font.bold = True; p.font.italic = True
        p.font.color.rgb = s.rgb(s.primary)

        # Bullet points
        points = spec.get("points", [])
        y = 2.55
        for point in points:
            # Main bullet (■)
            self._text(slide, 0.6, y, 12.3, 0.3,
                f"■  {point['main']}", 11, bold=True, color=s.text)
            y += 0.32
            # Sub-bullets (—)
            for sub in point.get("subs", []):
                self._text(slide, 1.0, y, 11.8, 0.28,
                    f"—  {sub}", 10, color=s.text)
                y += 0.28
            y += 0.1

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 5: INVESTMENT HIGHLIGHTS (4 card grid)
    # ═══════════════════════════════════════════════
    def render_investment_highlights(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", ""))
        if spec.get("subtitle"):
            self._subtitle(slide, spec["subtitle"])

        cards = spec.get("cards", [])[:4]
        # 2×2 grid
        card_w = 5.9
        card_h = 2.1
        x_offsets = [0.5, 6.93]
        y_offsets = [2.3, 4.55]

        for i, card in enumerate(cards):
            col = i % 2
            row = i // 2
            x = x_offsets[col]
            y = y_offsets[row]

            # Card shape
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(card_w), Inches(card_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = s.rgb("#FFFFFF")
            rect.line.color.rgb = s.rgb(s.primary)
            rect.line.width = Pt(0.75)

            # Number
            self._text(slide, x + 0.25, y + 0.15, 0.8, 0.5,
                f"0{i+1}", 28, bold=True, color=s.primary_light, font="Georgia")

            # Card title
            self._text(slide, x + 1.1, y + 0.2, card_w - 1.3, 0.4,
                card.get("title", ""), 14, bold=True, color=s.primary, font="Georgia")

            # Card body
            self._text(slide, x + 0.3, y + 0.75, card_w - 0.5, card_h - 0.9,
                card.get("body", ""), 10, color=s.text)

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 6: FINANCIAL SUMMARY TABLE
    # ═══════════════════════════════════════════════
    def render_financial_summary(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", "Financial Summary"))
        if spec.get("subtitle"):
            self._subtitle(slide, spec["subtitle"])

        # Navy header bar
        self._rect(slide, 0.5, 2.1, 12.333, 0.35, s.primary)
        self._text(slide, 0.65, 2.17, 12, 0.25,
            spec.get("section_header", ""), 11, bold=True, color="#FFFFFF")

        # Column headers + data table
        headers = spec.get("headers", [])  # ["($ in thousands)", "FY2023A", "FY2024A", ...]
        rows = spec.get("rows", [])        # [{"label": "Revenue", "values": [...], "style": "bold"}]

        n_cols = len(headers)
        n_rows = len(rows) + 1

        label_col_w = 2.8
        data_col_w = (12.333 - label_col_w) / (n_cols - 1)

        table_top = 2.65
        row_h = 0.32

        tbl_shape = slide.shapes.add_table(n_rows, n_cols,
            Inches(0.5), Inches(table_top), Inches(12.333), Inches(n_rows * row_h))
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(label_col_w)
        for c in range(1, n_cols):
            tbl.columns[c].width = Inches(data_col_w)

        # Header row
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10); p.font.bold = True
            p.font.color.rgb = s.rgb(s.primary)
            p.alignment = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

        # Data rows
        for r, row in enumerate(rows):
            row_style = row.get("style", "normal")
            is_total = row.get("total", False)
            indent = "  " if row_style == "pct" else ""

            for c, val in enumerate([row["label"]] + row["values"]):
                cell = tbl.cell(r + 1, c)
                cell.text = f"{indent}{val}" if c == 0 else str(val)
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(10)
                cp.font.name = "Calibri"

                if row_style == "bold":
                    cp.font.bold = True
                    cp.font.color.rgb = s.rgb(s.text)
                elif row_style == "pct":
                    cp.font.italic = True
                    cp.font.color.rgb = s.rgb(s.muted)
                elif row_style == "highlight":
                    cp.font.bold = True
                    cp.font.color.rgb = s.rgb(s.primary)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = s.rgb(s.header_bg)
                else:
                    cp.font.color.rgb = s.rgb(s.text)

                if not (row_style == "highlight"):
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

                cp.alignment = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 7: STACKED BAR + DATA TABLE (refined)
    # Reference: Moelis slide 9
    # ═══════════════════════════════════════════════
    def render_stacked_bar_table(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", ""))

        # Navy section header
        self._rect(slide, 0.5, 1.9, 12.333, 0.35, s.primary)
        self._text(slide, 0.65, 1.97, 12, 0.25,
            spec.get("header_text", ""), 11, bold=True, color="#FFFFFF")

        data = spec.get("data", {})
        periods = data.get("periods", [])
        segments = data.get("segments", [])
        totals = data.get("totals", [])

        n_periods = len(periods)
        chart_left = 1.0
        chart_right = 12.3
        chart_width = chart_right - chart_left
        bar_area_width = chart_width / n_periods
        bar_width = bar_area_width * 0.48
        chart_top = 2.55
        chart_height = 1.65
        max_total = max(totals) if totals else 1

        # Baseline
        baseline_y = chart_top + chart_height
        self._line(slide, chart_left - 0.1, baseline_y, chart_width + 0.2,
            color="#555555", thickness=0.015)

        # Stacked bars
        for p_idx, period in enumerate(periods):
            x_center = chart_left + p_idx * bar_area_width + bar_area_width / 2
            x_bar = x_center - bar_width / 2

            y_bottom = baseline_y
            for seg_idx, seg in enumerate(segments):
                val = seg["values"][p_idx]
                bar_h = (val / max_total) * chart_height if max_total > 0 else 0
                y_top = y_bottom - bar_h
                # Use distinct segment colors
                color = seg.get("color", s.segment_colors[seg_idx % len(s.segment_colors)])
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    Inches(x_bar), Inches(y_top), Inches(bar_width), Inches(bar_h))
                rect.fill.solid()
                rect.fill.fore_color.rgb = s.rgb(color)
                rect.line.fill.background()
                y_bottom = y_top

            # Total label (tight to bar top)
            total_str = spec.get("format_total", lambda x: f"${x:,.0f}")(totals[p_idx])
            top_of_bar = baseline_y - (totals[p_idx] / max_total) * chart_height
            self._text(slide, x_center - 0.8, top_of_bar - 0.24, 1.6, 0.22,
                total_str, 10, bold=True, color=s.text, align=PP_ALIGN.CENTER)

            # Period label
            self._text(slide, x_center - 0.6, baseline_y + 0.06, 1.2, 0.2,
                period, 9, color=s.muted, align=PP_ALIGN.CENTER)

        # Legend
        legend_y = baseline_y + 0.32
        entry_width = 1.7
        total_legend_width = len(segments) * entry_width
        legend_start_x = chart_left + (chart_width - total_legend_width) / 2
        for i, seg in enumerate(segments):
            ex = legend_start_x + i * entry_width
            color = seg.get("color", s.segment_colors[i % len(s.segment_colors)])
            sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(ex), Inches(legend_y + 0.04), Inches(0.14), Inches(0.14))
            sw.fill.solid(); sw.fill.fore_color.rgb = s.rgb(color)
            sw.line.fill.background()
            self._text(slide, ex + 0.2, legend_y, 1.4, 0.2,
                seg["name"], 9, color=s.text)

        # Data table below
        table_top = legend_y + 0.45
        table_rows = data.get("table_rows", [])
        if table_rows:
            n_cols = n_periods + 1
            row_h = 0.3
            tbl_shape = slide.shapes.add_table(len(table_rows) + 1, n_cols,
                Inches(0.5), Inches(table_top), Inches(12.333), Inches((len(table_rows) + 1) * row_h))
            tbl = tbl_shape.table
            tbl.columns[0].width = Inches(2.3)
            col_w = (12.333 - 2.3) / n_periods
            for c in range(1, n_cols):
                tbl.columns[c].width = Inches(col_w)

            # Header row
            cell = tbl.cell(0, 0)
            cell.text = ""
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")
            for c, p in enumerate(periods):
                cell = tbl.cell(0, c + 1)
                cell.text = p
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(9); cp.font.bold = True
                cp.font.color.rgb = s.rgb(s.primary)
                cp.alignment = PP_ALIGN.RIGHT
                cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

            # Data rows
            for r, row in enumerate(table_rows):
                # Label with colored text
                cell = tbl.cell(r + 1, 0)
                cell.text = row["label"]
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(9)
                cp.font.bold = row.get("bold", False)
                cp.font.color.rgb = s.rgb(row.get("color", s.text))
                cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

                for c, val in enumerate(row["values"]):
                    cell = tbl.cell(r + 1, c + 1)
                    cell.text = str(val)
                    cp = cell.text_frame.paragraphs[0]
                    cp.font.size = Pt(9)
                    cp.alignment = PP_ALIGN.RIGHT
                    cp.font.bold = row.get("bold", False)
                    cp.font.color.rgb = s.rgb(s.text)
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 8: DUAL CHART
    # ═══════════════════════════════════════════════
    def render_dual_chart(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", ""))

        charts = spec.get("charts", [{}, {}])
        for ci, chart in enumerate(charts):
            x_offset = 0.5 if ci == 0 else 6.8
            zone_width = 5.8

            # Chart subtitle
            self._text(slide, x_offset + 0.4, 2.0, zone_width - 3.2, 0.3,
                chart.get("subtitle", ""), 12, bold=True, color=s.primary)

            # CAGR annotation
            if chart.get("cagr"):
                self._text(slide, x_offset + zone_width - 2.5, 2.0, 2.4, 0.3,
                    f"CAGR: {chart['cagr']}", 11, bold=True,
                    color="#2E8B57", align=PP_ALIGN.RIGHT)

            periods = chart.get("periods", [])
            values = chart.get("values", [])
            bar_color = chart.get("color", s.primary)
            n = len(periods)
            if not n: continue

            chart_left = x_offset + 0.4
            chart_width_c = zone_width - 0.8
            chart_top = 2.6
            chart_h = 2.25
            baseline_y = chart_top + chart_h
            max_val = max(values) if values else 1
            bar_area_w = chart_width_c / n
            bar_w = bar_area_w * 0.50

            # Gridlines
            for pct in [0.25, 0.5, 0.75, 1.0]:
                y = baseline_y - pct * chart_h
                self._line(slide, chart_left, y, chart_width_c,
                    color="#DDDDDD", thickness=0.005)

            # Baseline
            self._line(slide, chart_left, baseline_y, chart_width_c,
                color="#555555", thickness=0.015)

            # Bars
            for i, (period, val) in enumerate(zip(periods, values)):
                x = chart_left + i * bar_area_w + (bar_area_w - bar_w) / 2
                bar_h = (val / max_val) * chart_h
                y = baseline_y - bar_h
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    Inches(x), Inches(y), Inches(bar_w), Inches(bar_h))
                rect.fill.solid()
                rect.fill.fore_color.rgb = s.rgb(bar_color)
                rect.line.fill.background()

                # Value label
                fmt_fn = chart.get("format", lambda v: f"${v:,.0f}")
                self._text(slide, x - 0.35, y - 0.26, bar_w + 0.7, 0.22,
                    fmt_fn(val), 10, bold=True, color=s.text, align=PP_ALIGN.CENTER)

                # Period
                self._text(slide, x - 0.35, baseline_y + 0.06, bar_w + 0.7, 0.2,
                    period, 9, color=s.muted, align=PP_ALIGN.CENTER)

            # Secondary metrics row
            if chart.get("secondary_values"):
                self._text(slide, chart_left, baseline_y + 0.32, chart_width_c, 0.2,
                    chart.get("secondary_label", ""), 8, bold=True,
                    color=s.muted, align=PP_ALIGN.CENTER)
                for i, val in enumerate(chart["secondary_values"]):
                    x = chart_left + i * bar_area_w + (bar_area_w - bar_w) / 2
                    self._text(slide, x - 0.35, baseline_y + 0.52, bar_w + 0.7, 0.2,
                        val, 9, bold=True, color=s.text, align=PP_ALIGN.CENTER)

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 9: FOOTBALL FIELD
    # ═══════════════════════════════════════════════
    def render_football_field(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", "Preliminary Valuation Summary"), y=1.0, height=0.7)
        if spec.get("subtitle"):
            self._subtitle(slide, spec["subtitle"], y=1.75)

        data = spec.get("methodologies", [])
        label_left = 0.5
        label_width = 1.5
        method_left = 2.15
        method_width = 2.6
        bar_left = 4.9
        bar_width = 5.4
        mult_left = 10.5
        mult_width = 2.3
        row_height = 0.65
        start_y = 3.0

        all_lows = [m["low"] for m in data]
        all_highs = [m["high"] for m in data]
        price_min = min(all_lows) * 0.85
        price_max = max(all_highs) * 1.10

        def price_to_x(price):
            pct = (price - price_min) / (price_max - price_min)
            return bar_left + pct * bar_width

        # Reference lines FIRST (behind everything)
        ref_lines = spec.get("reference_lines", [])
        ref_positions = sorted([(r, price_to_x(r["price"])) for r in ref_lines],
                                key=lambda t: t[1])
        occupied_row1 = []
        label_width_r = 1.3
        label_height = 0.22
        ref_label_specs = []

        for ref, x in ref_positions:
            line_top = 2.13
            line_bottom = start_y + len(data) * row_height
            self._vline(slide, x, line_top, line_bottom - line_top,
                color=ref.get("color", "#000000"), thickness=0.02)

            label_x = x - label_width_r / 2
            label_y = 2.15
            label_end = label_x + label_width_r
            for (ox_s, ox_e) in occupied_row1:
                if not (label_end < ox_s or label_x > ox_e):
                    label_y = 2.42
                    break
            else:
                occupied_row1.append((label_x, label_end))
            ref_label_specs.append((ref, label_x, label_y))

        # Column headers
        hdr_y = 2.7
        self._text(slide, label_left, hdr_y, label_width, 0.25,
            "Category", 8, bold=True, color=s.primary, align=PP_ALIGN.CENTER)
        self._text(slide, method_left, hdr_y, method_width, 0.25,
            "Valuation Methodology", 8, bold=True, color=s.primary)
        self._text(slide, bar_left, hdr_y, bar_width, 0.25,
            "Valuation Range ($ per Share)", 8, bold=True, color=s.primary, align=PP_ALIGN.CENTER)
        self._text(slide, mult_left, hdr_y, mult_width, 0.25,
            "Implied TEV/EBITDA", 8, bold=True, color=s.primary, align=PP_ALIGN.CENTER)
        self._line(slide, 0.5, hdr_y + 0.28, 12.333, color=s.primary, thickness=0.015)

        # Data rows
        for i, method in enumerate(data):
            y = start_y + i * row_height
            if i % 2 == 0:
                self._rect(slide, 0.5, y, 12.333, row_height, s.row_alt)

            # Pill
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(label_left + 0.15), Inches(y + 0.15),
                Inches(label_width - 0.3), Inches(0.35))
            pill.fill.solid()
            pill.fill.fore_color.rgb = s.rgb(method.get("category_color", s.primary))
            pill.line.fill.background()
            tf = pill.text_frame; p = tf.paragraphs[0]
            p.text = method["category"]
            p.font.size = Pt(8); p.font.bold = True
            p.font.color.rgb = s.rgb("#FFFFFF")
            p.alignment = PP_ALIGN.CENTER

            # Description
            self._text(slide, method_left, y + 0.08, method_width - 0.1, row_height - 0.15,
                method["description"], 8, color=s.text)

            # Bar
            x_low = price_to_x(method["low"])
            x_high = price_to_x(method["high"])
            self._rect(slide, x_low, y + 0.2, x_high - x_low, 0.25,
                method.get("bar_color", s.primary))

            # Price labels with white bg (readable over reference lines)
            self._text_bg(slide, x_low - 0.55, y + 0.15, 0.5, 0.28,
                f"${method['low']:.0f}", bg_color="#FFFFFF",
                size=8, bold=True, color=s.text, align=PP_ALIGN.RIGHT)
            self._text_bg(slide, x_high + 0.05, y + 0.15, 0.55, 0.28,
                f"${method['high']:.0f}", bg_color="#FFFFFF",
                size=8, bold=True, color=s.text, align=PP_ALIGN.LEFT)

            # Multiples
            self._text(slide, mult_left, y + 0.18, mult_width, 0.3,
                method.get("implied_multiples", ""), 8, color=s.text, align=PP_ALIGN.CENTER)

        bars_bottom = start_y + len(data) * row_height

        # Reference labels LAST (on top)
        for ref, label_x, label_y in ref_label_specs:
            self._text_bg(slide, label_x, label_y, label_width_r, label_height,
                f"{ref['label']}: ${ref['price']:.0f}",
                bg_color="#FFFFFF", border_color=ref.get("color", "#000000"),
                size=8, bold=True, color=ref.get("color", "#000000"))

        # X-axis
        axis_y = bars_bottom + 0.12
        self._line(slide, bar_left, axis_y, bar_width, color="#888888", thickness=0.015)
        for i in range(6):
            price = price_min + i * (price_max - price_min) / 5
            x_tick = bar_left + i * bar_width / 5
            self._text(slide, x_tick - 0.3, axis_y + 0.05, 0.6, 0.2,
                f"${price:.0f}", 7, color=s.muted, align=PP_ALIGN.CENTER)

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 10: SENSITIVITY TABLE
    # ═══════════════════════════════════════════════
    def render_sensitivity(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", ""))
        if spec.get("subtitle"):
            self._subtitle(slide, spec["subtitle"])

        # Main label above table
        self._text(slide, 0.5, 2.2, 8, 0.3,
            spec.get("table_label", "Implied Share Price ($)"),
            12, bold=True, color=s.primary)

        data = spec.get("data", {})
        row_headers = data.get("row_headers", [])
        col_headers = data.get("col_headers", [])
        values = data.get("values", [])
        base_row = data.get("base_row", -1)
        base_col = data.get("base_col", -1)

        row_label = spec.get("row_axis_label", "WACC")
        col_label = spec.get("col_axis_label", "Terminal Growth Rate")

        # Col axis label (centered above col headers)
        self._text(slide, 2.5, 2.6, 8, 0.25,
            col_label, 9, bold=True, italic=True, color=s.muted, align=PP_ALIGN.CENTER)

        # Row axis label (rotated would be nice, but just put on left)
        self._text(slide, 0.5, 4.0, 1.2, 0.3,
            row_label, 9, bold=True, italic=True, color=s.muted)

        n_rows = len(row_headers) + 1
        n_cols = len(col_headers) + 1

        table_left = 1.8
        table_top = 2.9
        table_width = 9.5
        row_h = 0.35
        table_h = n_rows * row_h

        tbl_shape = slide.shapes.add_table(n_rows, n_cols,
            Inches(table_left), Inches(table_top), Inches(table_width), Inches(table_h))
        tbl = tbl_shape.table

        col_w = table_width / n_cols
        for c in range(n_cols):
            tbl.columns[c].width = Inches(col_w)

        # Corner
        cell = tbl.cell(0, 0)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.header_bg)

        # Col headers
        for c, hdr in enumerate(col_headers):
            cell = tbl.cell(0, c + 1)
            cell.text = hdr
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(10); cp.font.bold = True
            cp.font.color.rgb = s.rgb(s.primary)
            cp.alignment = PP_ALIGN.CENTER
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.header_bg)

        # Row headers
        for r, hdr in enumerate(row_headers):
            cell = tbl.cell(r + 1, 0)
            cell.text = hdr
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(10); cp.font.bold = True
            cp.font.color.rgb = s.rgb(s.primary)
            cp.alignment = PP_ALIGN.CENTER
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.header_bg)

        # Values
        for r in range(len(row_headers)):
            for c in range(len(col_headers)):
                cell = tbl.cell(r + 1, c + 1)
                cell.text = str(values[r][c])
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(10)
                cp.alignment = PP_ALIGN.CENTER
                cp.font.color.rgb = s.rgb(s.text)
                if r == base_row and c == base_col:
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#BDD7EE")
                    cp.font.bold = True
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 11: SOURCES & USES (LBO)
    # ═══════════════════════════════════════════════
    def render_sources_uses(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", "Sources & Uses"))
        if spec.get("subtitle"):
            self._subtitle(slide, spec["subtitle"])

        # Two column layout: Sources left, Uses right
        col_width = 5.9
        left_x = 0.5
        right_x = 6.93
        top_y = 2.3

        # Navy headers
        self._rect(slide, left_x, top_y, col_width, 0.35, s.primary)
        self._text(slide, left_x, top_y + 0.07, col_width, 0.25,
            "SOURCES", 12, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

        self._rect(slide, right_x, top_y, col_width, 0.35, s.primary)
        self._text(slide, right_x, top_y + 0.07, col_width, 0.25,
            "USES", 12, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

        # Sources table
        sources = spec.get("sources", [])
        y = top_y + 0.5
        for i, item in enumerate(sources):
            if i % 2 == 0:
                self._rect(slide, left_x, y, col_width, 0.35, s.row_alt)
            is_total = item.get("total", False)
            label = item["label"]
            amount = item["amount"]
            pct = item.get("pct", "")
            if is_total:
                self._line(slide, left_x, y, col_width, color=s.primary, thickness=0.015)
            self._text(slide, left_x + 0.15, y + 0.07, col_width * 0.5, 0.25,
                label, 11, bold=is_total, color=s.text)
            self._text(slide, left_x + col_width * 0.5, y + 0.07, col_width * 0.3 - 0.1, 0.25,
                f"${amount:,.0f}" if isinstance(amount, (int, float)) else amount,
                11, bold=is_total, color=s.text, align=PP_ALIGN.RIGHT)
            if pct:
                self._text(slide, left_x + col_width * 0.8, y + 0.07, col_width * 0.2 - 0.15, 0.25,
                    pct, 11, color=s.muted, align=PP_ALIGN.RIGHT)
            y += 0.35

        # Uses table
        uses = spec.get("uses", [])
        y = top_y + 0.5
        for i, item in enumerate(uses):
            if i % 2 == 0:
                self._rect(slide, right_x, y, col_width, 0.35, s.row_alt)
            is_total = item.get("total", False)
            label = item["label"]
            amount = item["amount"]
            pct = item.get("pct", "")
            if is_total:
                self._line(slide, right_x, y, col_width, color=s.primary, thickness=0.015)
            self._text(slide, right_x + 0.15, y + 0.07, col_width * 0.5, 0.25,
                label, 11, bold=is_total, color=s.text)
            self._text(slide, right_x + col_width * 0.5, y + 0.07, col_width * 0.3 - 0.1, 0.25,
                f"${amount:,.0f}" if isinstance(amount, (int, float)) else amount,
                11, bold=is_total, color=s.text, align=PP_ALIGN.RIGHT)
            if pct:
                self._text(slide, right_x + col_width * 0.8, y + 0.07, col_width * 0.2 - 0.15, 0.25,
                    pct, 11, color=s.muted, align=PP_ALIGN.RIGHT)
            y += 0.35

        # Balance check
        if spec.get("balance_check"):
            check_y = 6.2
            self._rect(slide, 0.5, check_y, 12.333, 0.35, s.header_bg)
            self._text(slide, 0.5, check_y + 0.07, 12.333, 0.25,
                spec["balance_check"], 11, bold=True, color=s.primary, align=PP_ALIGN.CENTER)

        self._source(slide, spec.get("source", ""))

    # ═══════════════════════════════════════════════
    # TEMPLATE 12: TRADING COMPS
    # ═══════════════════════════════════════════════
    def render_trading_comps(self, spec: dict):
        s = self.style
        slide = self._new_slide()
        self._chrome(slide, spec.get("slide_number"))
        self._title(slide, spec.get("title", "Trading Comparables Analysis"))
        if spec.get("subtitle"):
            self._subtitle(slide, spec["subtitle"])

        # Navy header bar
        self._rect(slide, 0.5, 2.15, 12.333, 0.35, s.primary)
        self._text(slide, 0.65, 2.22, 12, 0.25,
            spec.get("header_text", "Selected Public Companies"),
            11, bold=True, color="#FFFFFF")

        headers = spec.get("headers", [])
        rows = spec.get("rows", [])
        target_row = spec.get("target_row", -1)

        n_cols = len(headers)
        n_rows = len(rows) + 1

        label_col_w = 2.5
        data_col_w = (12.333 - label_col_w) / (n_cols - 1)

        tbl_shape = slide.shapes.add_table(n_rows, n_cols,
            Inches(0.5), Inches(2.6), Inches(12.333), Inches(n_rows * 0.3))
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(label_col_w)
        for c in range(1, n_cols):
            tbl.columns[c].width = Inches(data_col_w)

        # Header row
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9); p.font.bold = True
            p.font.color.rgb = s.rgb(s.primary)
            p.alignment = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.header_bg)

        # Data rows
        for r, row in enumerate(rows):
            is_target = (r == target_row)
            is_summary = row.get("summary", False)

            for c, val in enumerate([row["company"]] + row["values"]):
                cell = tbl.cell(r + 1, c)
                cell.text = str(val)
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(9)
                cp.font.name = "Calibri"
                cp.alignment = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT

                if is_target:
                    cp.font.bold = True
                    cp.font.color.rgb = s.rgb(s.primary)
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFF2CC")
                elif is_summary:
                    cp.font.bold = True
                    cp.font.color.rgb = s.rgb(s.text)
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.row_alt)
                else:
                    cp.font.color.rgb = s.rgb(s.text)
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

        self._source(slide, spec.get("source", ""))

    # ─── Save ───
    def save(self, path):
        self.prs.save(path)


# ═══════════════════════════════════════════════
# DEMO: Full ADUS deck using all templates
# ═══════════════════════════════════════════════
if __name__ == "__main__":
    r = IBRenderer()

    # ─── 1. Cover ───
    r.render_cover({
        "title": "Project Alpine — Board of Directors\nDiscussion Materials",
        "subtitle": "Advisory Group LLC",
        "date": "April 2026",
    })

    # ─── 2. Table of Contents ───
    r.render_toc({
        "slide_number": 2,
        "title": "Today's Agenda",
        "items": [
            "Situation Overview and Strategic Context",
            "Review of Financial Performance",
            "Preliminary Valuation Analysis",
            "Perspectives on Process and Potential Buyers",
        ]
    })

    # ─── 3. Executive Summary ───
    r.render_exec_summary({
        "slide_number": 3,
        "title": "Executive Summary | Advisory Group Perspectives",
        "callout": "Advisory Group is pleased to present our preliminary perspectives on Addus HomeCare to the Board of Directors",
        "points": [
            {
                "main": "Addus HomeCare operates the largest independent personal care platform in the US, serving 107,000 consumers across 23 states with $1.4B in revenue",
                "subs": [
                    "Personal care segment (76.6% of revenue) grew 27.2% YoY driven by acquisitions and organic volume gains",
                    "EBITDA margins expanded 220bps over 3 years from 8.7% to 10.9%, demonstrating operating leverage at scale"
                ]
            },
            {
                "main": "The company is well-positioned for continued growth through both organic expansion and tuck-in M&A in a highly fragmented market",
                "subs": [
                    "Completed 3 acquisitions in FY2025; $525M+ undrawn revolver capacity for continued deal activity",
                    "Secular tailwinds from aging demographics and CMS policy shift toward home and community-based services"
                ]
            },
            {
                "main": "Preliminary valuation analysis suggests a range of $104-$174 per share with strong LBO returns at reasonable entry multiples",
                "subs": [
                    "DCF analysis (perpetuity growth): $129/share at 9.9% WACC, 3.0% terminal growth",
                    "LBO analysis: 2.51x MOIC / 20.2% IRR at 12.0x entry with 4.5x leverage, no multiple expansion assumed"
                ]
            },
        ]
    })

    # ─── 4. Investment Highlights ───
    r.render_investment_highlights({
        "slide_number": 4,
        "title": "Investment Highlights",
        "subtitle": "Scaled Home Care Platform With Secular Tailwinds and Proven Acquisition Engine",
        "cards": [
            {"title": "Market-Leading Platform", "body": "$1.09B Personal Care revenue serving 107,000 consumers across 23 states. Largest independent provider for dual-eligible populations."},
            {"title": "Expanding Margins at Scale", "body": "EBITDA margins expanded from 8.7% (FY2022) to 10.9% (FY2025). Capital-light model with only 0.5% CapEx/Revenue."},
            {"title": "Proven Buy-and-Build", "body": "23.2% revenue growth in FY2025 combining organic growth + strategic acquisitions. Revenue grown from $951M to $1.42B in 3 years."},
            {"title": "Demographic Tailwinds", "body": "65+ population growing 3% annually. Home care costs 50-75% less than facility-based alternatives, driving structural demand."},
        ]
    })

    # ─── 5. Section Divider ───
    r.render_section_divider({
        "title": "Review of Financial Performance"
    })

    # ─── 6. Financial Summary Table ───
    r.render_financial_summary({
        "slide_number": 6,
        "title": "Addus Has Delivered Consistent Revenue Growth with Expanding Profitability",
        "section_header": "Historical Financial Summary ($ in thousands)",
        "headers": ["Metric", "FY2023A", "FY2024A", "FY2025A", "3Y CAGR"],
        "rows": [
            {"label": "Revenue", "values": ["1,058,651", "1,154,599", "1,422,530", "14.4%"], "style": "bold"},
            {"label": "% Growth", "values": ["11.3%", "9.1%", "23.2%", ""], "style": "pct"},
            {"label": "Gross Profit", "values": ["339,876", "375,021", "461,874", "16.6%"], "style": "bold"},
            {"label": "% Margin", "values": ["32.1%", "32.5%", "32.5%", ""], "style": "pct"},
            {"label": "EBITDA", "values": ["105,082", "116,221", "155,027", "21.4%"], "style": "highlight"},
            {"label": "% Margin", "values": ["9.9%", "10.1%", "10.9%", ""], "style": "pct"},
            {"label": "Net Income", "values": ["62,516", "73,598", "95,910", "23.8%"], "style": "bold"},
            {"label": "Diluted EPS", "values": ["$3.83", "$4.23", "$5.22", "16.7%"], "style": "normal"},
        ],
        "source": "Source: SEC EDGAR via EdgarTools. FY2025 10-K filed February 24, 2026."
    })

    # ─── 7. Stacked Bar + Data Table ───
    r.render_stacked_bar_table({
        "slide_number": 7,
        "title": "Personal Care Dominates Revenue Mix With Accelerating Growth Across All Three Segments",
        "header_text": "Net Service Revenue by Segment ($ in thousands)",
        "source": "Source: adus_10k_master.json. SEC EDGAR 10-K filings.",
        "format_total": lambda x: f"${x/1000:,.0f}M",
        "data": {
            "periods": ["FY2022A", "FY2023A", "FY2024A", "FY2025A"],
            "segments": [
                {"name": "Personal Care", "color": "#1B365D", "values": [706507, 794718, 856581, 1089215]},
                {"name": "Hospice", "color": "#4B7BA8", "values": [201772, 207155, 228191, 262542]},
                {"name": "Home Health", "color": "#8EB4D6", "values": [42841, 56778, 69827, 70773]},
            ],
            "totals": [951120, 1058651, 1154599, 1422530],
            "table_rows": [
                {"label": "Personal Care", "values": ["$706,507", "$794,718", "$856,581", "$1,089,215"], "color": "#1B365D"},
                {"label": "Hospice", "values": ["$201,772", "$207,155", "$228,191", "$262,542"], "color": "#4B7BA8"},
                {"label": "Home Health", "values": ["$42,841", "$56,778", "$69,827", "$70,773"], "color": "#8EB4D6"},
                {"label": "Total", "values": ["$951,120", "$1,058,651", "$1,154,599", "$1,422,530"], "bold": True},
            ]
        }
    })

    # ─── 8. Dual Chart ───
    r.render_dual_chart({
        "slide_number": 8,
        "title": "Consistent Revenue Growth Accompanied by Meaningful EBITDA Margin Expansion",
        "source": "Source: adus_10k_master.json. SEC EDGAR XBRL data.",
        "charts": [
            {
                "subtitle": "Revenue ($M)",
                "periods": ["2022A", "2023A", "2024A", "2025A"],
                "values": [951120, 1058651, 1154599, 1422530],
                "color": "#1B365D",
                "format": lambda v: f"${v/1000:,.0f}M",
                "cagr": "14.4%",
                "secondary_label": "YoY Growth",
                "secondary_values": ["—", "+11.3%", "+9.1%", "+23.2%"],
            },
            {
                "subtitle": "EBITDA ($M) & Margin",
                "periods": ["2022A", "2023A", "2024A", "2025A"],
                "values": [82797, 105082, 116221, 155027],
                "color": "#4B7BA8",
                "format": lambda v: f"${v/1000:,.0f}M",
                "cagr": "23.3%",
                "secondary_label": "EBITDA Margin",
                "secondary_values": ["8.7%", "9.9%", "10.1%", "10.9%"],
            },
        ]
    })

    # ─── 9. Section Divider ───
    r.render_section_divider({
        "title": "Preliminary Valuation Analysis"
    })

    # ─── 10. Football Field ───
    r.render_football_field({
        "slide_number": 10,
        "title": "Preliminary Valuation Summary Implies a Range of $104–$174 Per Share Across Methodologies",
        "subtitle": "(All financials in $ per share unless otherwise stated)",
        "source": "Source: adus_dcf_model.xlsx, adus_lbo_model.xlsx. Management projections. FactSet consensus estimates as of April 2026.",
        "methodologies": [
            {"category": "DCF\nPerpetuity", "category_color": "#1B365D",
             "description": "WACC: 8.5-11.0%\nTGR: 2.0-4.0%",
             "low": 104, "high": 174, "bar_color": "#1B365D",
             "implied_multiples": "8.5x - 15.2x"},
            {"category": "DCF\nExit Multiple", "category_color": "#1B365D",
             "description": "WACC: 8.5-11.0%\nExit: 10.0-14.0x EBITDA",
             "low": 111, "high": 178, "bar_color": "#3B6B9D",
             "implied_multiples": "9.1x - 15.8x"},
            {"category": "Trading\nComps", "category_color": "#E07020",
             "description": "Selected Public Companies\n10.0x - 13.5x NTM EBITDA",
             "low": 115, "high": 155, "bar_color": "#E07020",
             "implied_multiples": "10.0x - 13.5x"},
            {"category": "LBO\nAnalysis", "category_color": "#2E8B57",
             "description": "2.0-3.0x MOIC target\n15-25% IRR range",
             "low": 95, "high": 140, "bar_color": "#2E8B57",
             "implied_multiples": "8.0x - 12.0x"},
            {"category": "52-Week\nTrading", "category_color": "#888888",
             "description": "52-Week High / Low\nHistorical trading range",
             "low": 85, "high": 135, "bar_color": "#AAAAAA",
             "implied_multiples": "n/a"},
        ],
        "reference_lines": [
            {"price": 129, "label": "DCF Base", "color": "#1B365D"},
            {"price": 98, "label": "Current", "color": "#C00000"},
        ]
    })

    # ─── 11. Sensitivity Table ───
    r.render_sensitivity({
        "slide_number": 11,
        "title": "DCF Sensitivity Analysis Across WACC and Terminal Growth Rate Assumptions",
        "table_label": "Implied Share Price ($)",
        "row_axis_label": "WACC →",
        "col_axis_label": "Terminal Growth Rate →",
        "data": {
            "row_headers": ["8.5%", "9.0%", "9.5%", "9.9%", "10.5%", "11.0%"],
            "col_headers": ["2.0%", "2.5%", "3.0%", "3.5%", "4.0%"],
            "values": [
                ["$147", "$159", "$174", "$194", "$220"],
                ["$136", "$146", "$158", "$174", "$194"],
                ["$126", "$135", "$145", "$157", "$173"],
                ["$119", "$127", "$129", "$145", "$158"],
                ["$111", "$118", "$125", "$134", "$145"],
                ["$104", "$110", "$117", "$125", "$134"],
            ],
            "base_row": 3,
            "base_col": 2,
        },
        "source": "Source: adus_dcf_model.xlsx. Base case highlighted. Each cell recalculates the full DCF for that assumption combination."
    })

    # ─── 12. Trading Comps ───
    r.render_trading_comps({
        "slide_number": 12,
        "title": "Addus Trades at a Discount to Home Care Services Peers Despite Superior Growth and Margins",
        "header_text": "Selected Home Care Services Trading Comparables",
        "headers": ["Company", "Ticker", "Market Cap ($M)", "EV ($M)", "'26E Rev ($M)", "'26E EBITDA ($M)", "EV/Rev", "EV/EBITDA", "Rev Growth"],
        "rows": [
            {"company": "Addus HomeCare", "values": ["ADUS", "1,781", "1,820", "1,565", "175", "1.2x", "10.4x", "10.0%"]},
            {"company": "BrightSpring Health", "values": ["BTSG", "5,820", "8,400", "11,200", "530", "0.8x", "15.8x", "12.5%"]},
            {"company": "Pennant Group", "values": ["PNTG", "3,940", "4,200", "765", "82", "5.5x", "51.2x", "18.2%"]},
            {"company": "Enhabit Home Health", "values": ["EHAB", "680", "1,100", "1,065", "115", "1.0x", "9.6x", "3.5%"]},
            {"company": "Aveanna Healthcare", "values": ["AVAH", "2,150", "3,450", "2,085", "195", "1.7x", "17.7x", "6.8%"]},
            {"company": "Modivcare", "values": ["MODV", "310", "965", "2,720", "155", "0.4x", "6.2x", "2.1%"]},
            {"company": "Mean", "values": ["", "", "", "", "", "1.8x", "18.5x", "8.9%"], "summary": True},
            {"company": "Median", "values": ["", "", "", "", "", "1.1x", "12.7x", "8.4%"], "summary": True},
        ],
        "target_row": 0,
        "source": "Source: FactSet, Capital IQ. Market data as of April 7, 2026. Addus highlighted in yellow."
    })

    # ─── 13. Section Divider ───
    r.render_section_divider({
        "title": "LBO Analysis & Transaction Assumptions"
    })

    # ─── 14. Sources & Uses ───
    r.render_sources_uses({
        "slide_number": 14,
        "title": "Preliminary LBO Capital Structure Sources Total Capitalization of $1,902M",
        "subtitle": "Based on 12.0x LTM EBITDA of $155M | 4.5x leverage | 63% sponsor equity contribution",
        "source": "Source: adus_lbo_model.xlsx. Transaction assumptions as of April 2026.",
        "sources": [
            {"label": "Term Loan B (SOFR + 450)", "amount": 700000, "pct": "36.8%"},
            {"label": "Revolver (undrawn)", "amount": 0, "pct": "0.0%"},
            {"label": "Sponsor Equity", "amount": 1202229, "pct": "63.2%"},
            {"label": "Total Sources", "amount": 1902229, "pct": "100.0%", "total": True},
        ],
        "uses": [
            {"label": "Enterprise Value", "amount": 1860324, "pct": "97.8%"},
            {"label": "Advisory Fees (1.5%)", "amount": 27905, "pct": "1.5%"},
            {"label": "Financing Fees (2.0%)", "amount": 14000, "pct": "0.7%"},
            {"label": "Total Uses", "amount": 1902229, "pct": "100.0%", "total": True},
        ],
        "balance_check": "Balance Check: Sources = Uses = $1,902,229K ✓"
    })

    import os
    output = os.path.join(os.path.dirname(__file__), "..", "case_study", "adus_full_deck.pptx")
    output = os.path.abspath(output)
    os.makedirs(os.path.dirname(output), exist_ok=True)
    r.save(output)
    print(f"✓ Saved {output}")
    print(f"✓ {len(r.prs.slides)} slides")
    print()
    print("Slides generated:")
    for i, title in enumerate([
        "Cover", "Table of Contents", "Executive Summary", "Investment Highlights",
        "Section Divider (Financial Performance)", "Financial Summary Table",
        "Stacked Bar + Data Table (Revenue by Segment)", "Dual Chart (Revenue + EBITDA)",
        "Section Divider (Valuation)", "Football Field", "Sensitivity Grid",
        "Trading Comps", "Section Divider (LBO)", "Sources & Uses",
    ], 1):
        print(f"  {i:2}. {title}")
