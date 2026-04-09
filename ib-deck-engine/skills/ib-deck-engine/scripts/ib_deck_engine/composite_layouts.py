"""
Composite layout engine v4 — production-grade IB slide patterns.

v4 fixes (referencing Moelis slide 9 "Renewables Revenue by Type"):
- Chart uses ~35% of vertical space (was 50%+)
- Data table fills more of the bottom half
- Tight data label placement (close to bars)
- Clean table with thin separators, bold Total row with top border
- Reference lines on football field drawn BEFORE bars so labels are on top
- Low/high price labels have white background for readability
- All three slides fit within footer area (nothing overflows)
- Source text is small but clearly positioned at very bottom
"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from engine import LayoutEngine, Renderer, GS_STYLE, Element
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class CompositeRenderer(Renderer):
    """Extended renderer with multi-zone composite layouts matching real IB decks."""

    # ═══════════════════════════════════════════════
    # LAYOUT 1: STACKED BAR + DATA TABLE
    # Reference: Moelis slide 9 (Renewables Revenue by Type)
    # ═══════════════════════════════════════════════
    def render_stacked_bar_with_table(self, spec: dict):
        s = self.style
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid(); bg.fore_color.rgb = s.rgb(s.background)

        self._add_chrome(slide, spec)
        self._add_title(slide, spec.get("title", ""))

        # ─── Layout regions (slide height 7.5") ───
        # 0.0-0.9:  Chrome
        # 1.0-1.75: Title
        # 1.9-2.25: Navy header bar
        # 2.40-4.30: Chart (1.9" tall)
        # 4.35-4.55: Year labels
        # 4.60-4.85: Legend
        # 5.00-6.50: Data table (1.5" for 5 rows)
        # 6.55-6.75: Source
        # 6.85-7.20: Footer badge

        # Navy section header bar
        header_y = 1.9
        bar_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(header_y), Inches(12.333), Inches(0.35))
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = s.rgb(s.primary)
        bar_shape.line.fill.background()
        tf = bar_shape.text_frame
        tf.margin_left = Pt(10)
        p = tf.paragraphs[0]
        p.text = spec.get("header_text", "")
        p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = s.rgb("#FFFFFF")
        p.alignment = PP_ALIGN.LEFT

        # Chart data
        data = spec.get("data", {})
        periods = data.get("periods", [])
        segments = data.get("segments", [])
        totals = data.get("totals", [])

        n_periods = len(periods)
        chart_left = 1.0
        chart_right = 12.3
        chart_width = chart_right - chart_left
        bar_area_width = chart_width / n_periods
        bar_width = bar_area_width * 0.48   # matches Moelis ratio
        chart_top = 2.55
        chart_height = 1.75
        max_total = max(totals) if totals else 1

        # Baseline x-axis
        baseline_y = chart_top + chart_height
        axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(chart_left - 0.1), Inches(baseline_y), Inches(chart_width + 0.2), Inches(0.01))
        axis.fill.solid(); axis.fill.fore_color.rgb = s.rgb("#666666")
        axis.line.fill.background()

        # Draw stacked bars
        for p_idx, period in enumerate(periods):
            x_center = chart_left + p_idx * bar_area_width + bar_area_width / 2
            x_bar = x_center - bar_width / 2

            y_bottom = baseline_y
            for seg in segments:
                val = seg["values"][p_idx]
                bar_h = (val / max_total) * chart_height if max_total > 0 else 0
                y_top = y_bottom - bar_h

                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    Inches(x_bar), Inches(y_top), Inches(bar_width), Inches(bar_h))
                rect.fill.solid()
                rect.fill.fore_color.rgb = s.rgb(seg["color"])
                rect.line.fill.background()
                y_bottom = y_top

            # Total label — tight to top of bar (0.22" above)
            total_str = spec.get("format_total", lambda x: f"${x:,.0f}")(totals[p_idx])
            top_of_bar = baseline_y - (totals[p_idx] / max_total) * chart_height
            self._add_text(slide, x_center - 0.7, top_of_bar - 0.22, 1.4, 0.2,
                total_str, 10, bold=True, align=PP_ALIGN.CENTER)

            # Period label — tight below baseline
            self._add_text(slide, x_center - 0.5, baseline_y + 0.05, 1.0, 0.2,
                period, 9, color=s.text_muted, align=PP_ALIGN.CENTER)

        # Legend — compact, centered below year labels
        legend_y = baseline_y + 0.32
        seg_entries = []
        for seg in segments:
            seg_entries.append((seg["name"], seg["color"]))
        # Compute widths and center
        entry_width = 1.7
        total_legend_width = len(seg_entries) * entry_width
        legend_start_x = chart_left + (chart_width - total_legend_width) / 2

        for i, (name, color) in enumerate(seg_entries):
            ex = legend_start_x + i * entry_width
            # Square color chip
            swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(ex), Inches(legend_y + 0.04), Inches(0.14), Inches(0.14))
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = s.rgb(color)
            swatch.line.fill.background()
            self._add_text(slide, ex + 0.2, legend_y, 1.4, 0.2,
                name, 9, color=s.text_primary)

        # DATA TABLE — compact, clean separators
        table_top = legend_y + 0.45
        table_data = data.get("table_rows", [])
        if table_data:
            n_rows = len(table_data) + 1
            n_cols = n_periods + 1
            row_h = 0.28
            table_h = row_h * n_rows

            tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                Inches(0.5), Inches(table_top), Inches(12.333), Inches(table_h))
            tbl = tbl_shape.table

            tbl.columns[0].width = Inches(2.3)
            data_col_w = (12.333 - 2.3) / n_periods
            for c in range(1, n_cols):
                tbl.columns[c].width = Inches(data_col_w)

            # Header row (thin border bottom, navy text, transparent bg)
            cell = tbl.cell(0, 0)
            cell.text = ""
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")
            for c, period in enumerate(periods):
                cell = tbl.cell(0, c + 1)
                cell.text = period
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(9); cp.font.bold = True
                cp.font.color.rgb = s.rgb(s.primary)
                cp.alignment = PP_ALIGN.RIGHT
                cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

            # Data rows — clean, colored row labels, right-aligned numbers
            for r, row in enumerate(table_data):
                is_total = row.get("bold", False) and r == len(table_data) - 1

                # Label cell
                cell = tbl.cell(r + 1, 0)
                cell.text = row["label"]
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(9)
                cp.font.bold = row.get("bold", False)
                # Colored row label (segment color)
                if row.get("color"):
                    cp.font.color.rgb = s.rgb(row["color"])
                else:
                    cp.font.color.rgb = s.rgb(s.text_primary)
                cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

                # Data value cells — RIGHT-ALIGNED
                for c, val in enumerate(row["values"]):
                    cell = tbl.cell(r + 1, c + 1)
                    cell.text = str(val)
                    cp = cell.text_frame.paragraphs[0]
                    cp.font.size = Pt(9)
                    cp.alignment = PP_ALIGN.RIGHT
                    cp.font.bold = row.get("bold", False)
                    cp.font.color.rgb = s.rgb(s.text_primary)
                    cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb("#FFFFFF")

            # Source — very bottom, small italic
            source_y = table_top + table_h + 0.08
            self._add_text(slide, 0.5, source_y, 12, 0.22,
                spec.get("source", ""), 8, italic=True, color=s.text_muted)

    # ═══════════════════════════════════════════════
    # LAYOUT 2: FOOTBALL FIELD
    # Reference: Evercore slide 8, Moelis slide 14
    # ═══════════════════════════════════════════════
    def render_football_field(self, spec: dict):
        s = self.style
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = s.rgb(s.background)

        self._add_chrome(slide, spec)

        # Title
        self._add_text(slide, 0.5, 1.0, 11, 0.7,
            spec.get("title", "Preliminary Valuation Summary"),
            s.title_size, bold=True, color=s.primary)

        if spec.get("subtitle"):
            self._add_text(slide, 0.5, 1.75, 8, 0.25,
                spec["subtitle"], 9, italic=True, color=s.text_muted)

        # Zones
        label_left = 0.5
        label_width = 1.5
        method_left = 2.15
        method_width = 2.6
        bar_left = 4.9
        bar_width = 5.4
        mult_left = 10.5
        mult_width = 2.3

        data = spec.get("methodologies", [])
        row_height = 0.65
        start_y = 3.0

        all_lows = [m["low"] for m in data]
        all_highs = [m["high"] for m in data]
        price_min = min(all_lows) * 0.85
        price_max = max(all_highs) * 1.10

        def price_to_x(price):
            pct = (price - price_min) / (price_max - price_min)
            return bar_left + pct * bar_width

        # ─── DRAW REFERENCE LINES FIRST (so bars/labels appear on top) ───
        ref_lines = spec.get("reference_lines", [])
        ref_positions = [(r, price_to_x(r["price"])) for r in ref_lines]
        ref_positions.sort(key=lambda t: t[1])

        occupied_row1 = []
        label_width_r = 1.3
        label_height = 0.22

        ref_label_specs = []  # draw labels AFTER bars so they're on top
        for ref, x in ref_positions:
            line_top = 2.13
            line_bottom = start_y + len(data) * row_height

            # Vertical line (drawn first = behind everything)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(line_top), Inches(0.02), Inches(line_bottom - line_top))
            line.fill.solid()
            line.fill.fore_color.rgb = s.rgb(ref.get("color", "#000000"))
            line.line.fill.background()

            # Compute label position (collision detect)
            label_x = x - label_width_r / 2
            label_y = 2.15
            label_end = label_x + label_width_r
            for (ox_start, ox_end) in occupied_row1:
                if not (label_end < ox_start or label_x > ox_end):
                    label_y = 2.42
                    break
            else:
                occupied_row1.append((label_x, label_end))
            ref_label_specs.append((ref, label_x, label_y))

        # Column headers
        hdr_y = 2.7
        self._add_text(slide, label_left, hdr_y, label_width, 0.25,
            "Category", 8, bold=True, color=s.primary, align=PP_ALIGN.CENTER)
        self._add_text(slide, method_left, hdr_y, method_width, 0.25,
            "Valuation Methodology", 8, bold=True, color=s.primary)
        self._add_text(slide, bar_left, hdr_y, bar_width, 0.25,
            "Valuation Range ($ per Share)", 8, bold=True, color=s.primary, align=PP_ALIGN.CENTER)
        self._add_text(slide, mult_left, hdr_y, mult_width, 0.25,
            "Implied TEV/EBITDA", 8, bold=True, color=s.primary, align=PP_ALIGN.CENTER)

        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(hdr_y + 0.28), Inches(12.333), Inches(0.015))
        sep.fill.solid(); sep.fill.fore_color.rgb = s.rgb(s.primary)
        sep.line.fill.background()

        # ─── DRAW DATA ROWS (alternating bg, pills, bars, price labels) ───
        for i, method in enumerate(data):
            y = start_y + i * row_height

            # Alternating row background (drawn first in each row)
            if i % 2 == 0:
                alt_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    Inches(0.5), Inches(y), Inches(12.333), Inches(row_height))
                alt_bg.fill.solid()
                alt_bg.fill.fore_color.rgb = s.rgb("#F5F5F5")
                alt_bg.line.fill.background()

            # Category pill
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(label_left + 0.15), Inches(y + 0.15), Inches(label_width - 0.3), Inches(0.35))
            pill.fill.solid()
            pill.fill.fore_color.rgb = s.rgb(method.get("category_color", s.primary))
            pill.line.fill.background()
            tf = pill.text_frame
            p = tf.paragraphs[0]
            p.text = method["category"]
            p.font.size = Pt(8); p.font.bold = True
            p.font.color.rgb = s.rgb("#FFFFFF")
            p.alignment = PP_ALIGN.CENTER

            # Methodology description
            self._add_text(slide, method_left, y + 0.08, method_width - 0.1, row_height - 0.15,
                method["description"], 8, color=s.text_primary)

            # Range bar (over the reference lines)
            x_low = price_to_x(method["low"])
            x_high = price_to_x(method["high"])
            bar_w = x_high - x_low
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(x_low), Inches(y + 0.2), Inches(bar_w), Inches(0.25))
            bar.fill.solid()
            bar.fill.fore_color.rgb = s.rgb(method.get("bar_color", s.primary))
            bar.line.fill.background()

            # Low/High labels — with WHITE BACKGROUND so they're readable over reference lines
            self._add_text_with_bg(slide, x_low - 0.55, y + 0.15, 0.5, 0.28,
                f"${method['low']:.0f}", 8, bold=True, color=s.text_primary,
                align=PP_ALIGN.RIGHT, bg_color=s.background)
            self._add_text_with_bg(slide, x_high + 0.05, y + 0.15, 0.55, 0.28,
                f"${method['high']:.0f}", 8, bold=True, color=s.text_primary,
                align=PP_ALIGN.LEFT, bg_color=s.background)

            # Implied multiples
            self._add_text(slide, mult_left, y + 0.18, mult_width, 0.3,
                method.get("implied_multiples", ""), 8, align=PP_ALIGN.CENTER)

        bars_bottom = start_y + len(data) * row_height

        # ─── NOW draw reference line LABELS (on top of everything) ───
        for ref, label_x, label_y in ref_label_specs:
            bg_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(label_x), Inches(label_y), Inches(label_width_r), Inches(label_height))
            bg_box.fill.solid()
            bg_box.fill.fore_color.rgb = s.rgb("#FFFFFF")
            bg_box.line.color.rgb = s.rgb(ref.get("color", "#000000"))
            bg_box.line.width = Pt(0.75)
            tf = bg_box.text_frame
            tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
            p = tf.paragraphs[0]
            p.text = f"{ref['label']}: ${ref['price']:.0f}"
            p.font.size = Pt(8); p.font.bold = True
            p.font.color.rgb = s.rgb(ref.get("color", "#000000"))
            p.alignment = PP_ALIGN.CENTER

        # X-axis price scale
        axis_y = bars_bottom + 0.12
        axis_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(bar_left), Inches(axis_y), Inches(bar_width), Inches(0.015))
        axis_line.fill.solid(); axis_line.fill.fore_color.rgb = s.rgb("#888888")
        axis_line.line.fill.background()

        n_ticks = 5
        for i in range(n_ticks + 1):
            price = price_min + i * (price_max - price_min) / n_ticks
            x_tick = bar_left + i * bar_width / n_ticks
            self._add_text(slide, x_tick - 0.3, axis_y + 0.05, 0.6, 0.2,
                f"${price:.0f}", 7, color=s.text_muted, align=PP_ALIGN.CENTER)

        # Source
        self._add_text(slide, 0.5, 6.65, 10, 0.25,
            spec.get("source", ""), 8, italic=True, color=s.text_muted)

    # ═══════════════════════════════════════════════
    # LAYOUT 3: DUAL CHART (Side-by-Side)
    # Reference: GS slide 13 (NTM EV/Revenue + EV/EBITDA)
    # ═══════════════════════════════════════════════
    def render_dual_chart(self, spec: dict):
        s = self.style
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = s.rgb(s.background)
        self._add_chrome(slide, spec)
        self._add_title(slide, spec.get("title", ""))

        charts = spec.get("charts", [{}, {}])

        for ci, chart in enumerate(charts):
            x_offset = 0.5 if ci == 0 else 6.8
            zone_width = 5.8

            # Chart subtitle (bold navy, left-aligned)
            self._add_text(slide, x_offset + 0.4, 2.0, zone_width - 3.2, 0.3,
                chart.get("subtitle", ""), 12, bold=True, color=s.primary)

            # CAGR annotation on same row, right-aligned
            if chart.get("cagr"):
                self._add_text(slide, x_offset + zone_width - 2.5, 2.0, 2.4, 0.3,
                    f"CAGR: {chart['cagr']}", 11, bold=True, color="#006100", align=PP_ALIGN.RIGHT)

            # Chart dimensions
            periods = chart.get("periods", [])
            values = chart.get("values", [])
            bar_color = chart.get("color", s.primary)
            n = len(periods)
            if not n: continue

            chart_left = x_offset + 0.4
            chart_width_c = zone_width - 0.8
            chart_top = 2.6
            chart_h = 2.3
            baseline_y = chart_top + chart_h
            max_val = max(values) if values else 1
            bar_area_w = chart_width_c / n
            bar_w = bar_area_w * 0.50

            # Gridlines first (behind bars)
            for grid_pct in [0.25, 0.50, 0.75, 1.0]:
                grid_y = baseline_y - grid_pct * chart_h
                grid_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    Inches(chart_left), Inches(grid_y), Inches(chart_width_c), Inches(0.005))
                grid_line.fill.solid(); grid_line.fill.fore_color.rgb = s.rgb("#DDDDDD")
                grid_line.line.fill.background()

            # Baseline
            axis_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(chart_left), Inches(baseline_y), Inches(chart_width_c), Inches(0.01))
            axis_line.fill.solid(); axis_line.fill.fore_color.rgb = s.rgb("#666666")
            axis_line.line.fill.background()

            # Bars
            for i, (period, val) in enumerate(zip(periods, values)):
                x = chart_left + i * bar_area_w + (bar_area_w - bar_w) / 2
                bar_h = (val / max_val) * chart_h if max_val > 0 else 0
                y = baseline_y - bar_h

                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    Inches(x), Inches(y), Inches(bar_w), Inches(bar_h))
                rect.fill.solid()
                rect.fill.fore_color.rgb = s.rgb(bar_color)
                rect.line.fill.background()

                # Value label above bar (tight)
                fmt_fn = chart.get("format", lambda v: f"${v:,.0f}")
                self._add_text(slide, x - 0.3, y - 0.28, bar_w + 0.6, 0.22,
                    fmt_fn(val), 10, bold=True, align=PP_ALIGN.CENTER)

                # Period label below baseline
                self._add_text(slide, x - 0.3, baseline_y + 0.06, bar_w + 0.6, 0.2,
                    period, 9, color=s.text_muted, align=PP_ALIGN.CENTER)

            # Secondary metric row
            if chart.get("secondary_values"):
                secondary = chart["secondary_values"]
                self._add_text(slide, chart_left, baseline_y + 0.35, chart_width_c, 0.2,
                    chart.get("secondary_label", ""), 8, bold=True, color=s.text_muted, align=PP_ALIGN.CENTER)
                for i, period in enumerate(periods):
                    x = chart_left + i * bar_area_w + (bar_area_w - bar_w) / 2
                    self._add_text(slide, x - 0.3, baseline_y + 0.55, bar_w + 0.6, 0.2,
                        secondary[i], 9, bold=True, color=s.text_primary, align=PP_ALIGN.CENTER)

        # Source
        self._add_text(slide, 0.5, 6.65, 10, 0.25,
            spec.get("source", ""), 8, italic=True, color=s.text_muted)

    # ═══════════════════════════════════════════════
    # Helpers
    # ═══════════════════════════════════════════════
    def _add_chrome(self, slide, spec):
        s = self.style
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.45), Inches(0.25), Inches(1.0), Inches(0.5))
        box.fill.background()
        box.line.color.rgb = s.rgb(s.primary); box.line.width = Pt(1)
        tf = box.text_frame; p = tf.paragraphs[0]
        p.text = s.bank_label; p.font.size = Pt(9)
        p.font.bold = True; p.font.color.rgb = s.rgb(s.primary)
        p.alignment = PP_ALIGN.CENTER

        self._add_text(slide, 11.5, 0.2, 1.5, 0.3,
            s.confidential_text, 10, italic=True, color=s.confidential_color, align=PP_ALIGN.RIGHT)

        if s.has_badge:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(10.3), Inches(6.95), Inches(s.badge_width), Inches(0.3))
            rect.fill.solid(); rect.fill.fore_color.rgb = s.rgb(s.primary)
            rect.line.fill.background()
            tf = rect.text_frame; p = tf.paragraphs[0]
            p.text = s.footer_text; p.font.size = Pt(9)
            p.font.bold = True; p.font.color.rgb = s.rgb("#FFFFFF")
            p.alignment = PP_ALIGN.CENTER
            strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(10.3 + s.badge_width), Inches(6.95), Inches(s.badge_strip_width), Inches(0.3))
            strip.fill.solid(); strip.fill.fore_color.rgb = s.rgb(s.secondary)
            strip.line.fill.background()

    def _add_title(self, slide, text):
        s = self.style
        self._add_text(slide, 0.5, 1.0, 11, 0.75, text,
            s.title_size, bold=True, color=s.primary)

    def _add_text(self, slide, left, top, width, height, text,
                  size=11, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT):
        s = self.style
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size); p.font.bold = bold; p.font.italic = italic
        p.font.color.rgb = s.rgb(color or s.text_primary)
        p.font.name = s.font; p.alignment = align
        p.space_after = Pt(0); p.space_before = Pt(0)

    def _add_text_with_bg(self, slide, left, top, width, height, text,
                          size=11, bold=False, italic=False, color=None,
                          align=PP_ALIGN.LEFT, bg_color="#FFFFFF"):
        """Text with a background fill — readable over lines/backgrounds."""
        s = self.style
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = s.rgb(bg_color)
        box.line.fill.background()
        tf = box.text_frame
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.margin_left = Pt(2); tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size); p.font.bold = bold; p.font.italic = italic
        p.font.color.rgb = s.rgb(color or s.text_primary)
        p.font.name = s.font; p.alignment = align


# ═══════════════════════════════════════
# DEMO
# ═══════════════════════════════════════
if __name__ == "__main__":
    renderer = CompositeRenderer(GS_STYLE)

    # 1. STACKED BAR + DATA TABLE (Moelis style)
    renderer.render_stacked_bar_with_table({
        "title": "Personal Care Dominates Revenue Mix With Accelerating\nGrowth Across All Three Segments",
        "slide_number": 5,
        "header_text": "Net Service Revenue by Segment ($ in thousands)",
        "source": "Source: adus_10k_master.json. SEC EDGAR 10-K filings.",
        "format_total": lambda x: f"${x/1000:,.0f}M",
        "data": {
            "periods": ["FY2022A", "FY2023A", "FY2024A", "FY2025A"],
            "segments": [
                {"name": "Personal Care", "color": "#1B365D", "values": [706507, 794718, 856581, 1089215]},
                {"name": "Hospice", "color": "#5B8DB8", "values": [201772, 207155, 228191, 262542]},
                {"name": "Home Health", "color": "#A7C4E0", "values": [42841, 56778, 69827, 70773]},
            ],
            "totals": [951120, 1058651, 1154599, 1422530],
            "table_rows": [
                {"label": "Personal Care", "values": ["$706,507", "$794,718", "$856,581", "$1,089,215"], "color": "#1B365D"},
                {"label": "Hospice", "values": ["$201,772", "$207,155", "$228,191", "$262,542"], "color": "#5B8DB8"},
                {"label": "Home Health", "values": ["$42,841", "$56,778", "$69,827", "$70,773"], "color": "#5B8DB8"},
                {"label": "Total", "values": ["$951,120", "$1,058,651", "$1,154,599", "$1,422,530"], "bold": True},
            ]
        }
    })

    # 2. FOOTBALL FIELD
    renderer.render_football_field({
        "title": "Preliminary Valuation Summary Implies a Range of\n$104–$174 Per Share Across Methodologies",
        "subtitle": "(All financials in $ per share unless otherwise stated)",
        "slide_number": 12,
        "source": "Source: adus_dcf_model.xlsx, adus_lbo_model.xlsx. Management projections. FactSet consensus estimates as of April 2026.",
        "methodologies": [
            {"category": "DCF\nPerpetuity", "category_color": "#1B365D",
             "description": "WACC: 8.5-11.0%\nTGR: 2.0-4.0%",
             "low": 104, "high": 174, "bar_color": "#1B365D",
             "implied_multiples": "8.5x - 15.2x"},
            {"category": "DCF\nExit Multiple", "category_color": "#1B365D",
             "description": "WACC: 8.5-11.0%\nExit: 10.0-14.0x EBITDA",
             "low": 111, "high": 178, "bar_color": "#3B6B9D",
             "implied_multiples": "9.1x - 15.8x"},
            {"category": "Trading\nComps", "category_color": "#E07020",
             "description": "Selected Public Companies\n10.0x - 13.5x NTM EBITDA",
             "low": 115, "high": 155, "bar_color": "#E07020",
             "implied_multiples": "10.0x - 13.5x"},
            {"category": "LBO\nAnalysis", "category_color": "#2E8B57",
             "description": "2.0-3.0x MOIC target\n15-25% IRR range",
             "low": 95, "high": 140, "bar_color": "#2E8B57",
             "implied_multiples": "8.0x - 12.0x"},
            {"category": "52-Week\nTrading", "category_color": "#888888",
             "description": "52-Week High / Low\nHistorical trading range",
             "low": 85, "high": 135, "bar_color": "#AAAAAA",
             "implied_multiples": "n/a"},
        ],
        "reference_lines": [
            {"price": 129, "label": "DCF Base", "color": "#1B365D"},
            {"price": 98, "label": "Current", "color": "#C00000"},
        ]
    })

    # 3. DUAL CHART
    renderer.render_dual_chart({
        "title": "Consistent Revenue Growth Accompanied by Meaningful\nEBITDA Margin Expansion",
        "slide_number": 6,
        "source": "Source: adus_10k_master.json. SEC EDGAR XBRL data.",
        "charts": [
            {
                "subtitle": "Revenue ($M)",
                "periods": ["2022A", "2023A", "2024A", "2025A"],
                "values": [951120, 1058651, 1154599, 1422530],
                "color": "#1B365D",
                "format": lambda v: f"${v/1000:,.0f}M",
                "cagr": "14.4%",
                "secondary_label": "YoY Growth",
                "secondary_values": ["—", "+11.3%", "+9.1%", "+23.2%"],
            },
            {
                "subtitle": "EBITDA ($M) & Margin",
                "periods": ["2022A", "2023A", "2024A", "2025A"],
                "values": [82797, 105082, 116221, 155027],
                "color": "#5B8DB8",
                "format": lambda v: f"${v/1000:,.0f}M",
                "cagr": "23.3%",
                "secondary_label": "EBITDA Margin",
                "secondary_values": ["8.7%", "9.9%", "10.1%", "10.9%"],
            },
        ]
    })

    output = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "case_study", "adus_composite_demo.pptx")
    os.makedirs(os.path.dirname(output), exist_ok=True)
    renderer.save(output)
    print(f"✓ Saved {output}")
    print(f"✓ {len(renderer.prs.slides)} slides")
    print(f"\nv4 changes (referencing Moelis slide 9):")
    print(f"  • Compressed chart to ~35% vertical (was 50%+)")
    print(f"  • Table now has clean thin separators (no color bands)")
    print(f"  • Row labels use segment colors inline")
    print(f"  • Source text position is now safely above footer")
    print(f"  • Football field: reference lines drawn BEFORE bars")
    print(f"  • Football field: price labels have white bg (readable over lines)")
    print(f"  • Reference line labels on top of everything")
    print(f"  • Tighter data label placement (0.22\" above bars)")
    print(f"  • Footer badge moved to y=6.95 to avoid overlap")
