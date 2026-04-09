"""
IB Deck Engine — A layout engine for investment banking presentations.
Inspired by Beautiful.ai's architecture, built for IB-specific patterns.

Architecture:
  Content (JSON) → Layout Engine (design rules) → Renderer (python-pptx) → PPTX

The LLM produces the JSON. The engine handles every pixel.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from dataclasses import dataclass, field
from typing import Optional


# ═══════════════════════════════════════════════
# DESIGN TOKENS — Configurable per bank
# ═══════════════════════════════════════════════

@dataclass
class BankStyle:
    """Design tokens for a specific bank's presentation style."""
    name: str
    # Colors
    primary: str = "#1B365D"        # Main brand color (navy for GS/Moelis)
    secondary: str = "#A7C4E0"      # Accent (light blue)
    background: str = "#EDEDED"     # Slide background
    confidential_color: str = "#C00000"  # Red
    text_primary: str = "#000000"
    text_muted: str = "#888888"
    header_fill: str = "#D6E4F0"    # Light blue for table headers
    highlight: str = "#BDD7EE"      # Base case highlight
    alt_row: str = "#F2F2F2"        # Alternating row
    positive: str = "#006100"       # Green for positive
    negative: str = "#C00000"       # Red for negative
    # Typography
    font: str = "Calibri"
    title_size: int = 24
    body_size: int = 11
    table_size: int = 10
    footer_size: int = 8
    # Layout
    margin_left: float = 0.5        # inches
    margin_right: float = 0.5
    margin_top: float = 1.2         # below logo/confidential
    margin_bottom: float = 0.8
    content_width: float = 12.333   # 13.333 - margins
    # Footer
    footer_text: str = "Investment Banking"
    confidential_text: str = "Confidential"
    bank_label: str = "Advisory Group"
    # Badge
    has_badge: bool = True
    badge_width: float = 2.2
    badge_strip_width: float = 0.4

    def rgb(self, hex_color: str) -> RGBColor:
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Pre-built bank styles
GS_STYLE = BankStyle(
    name="Goldman Sachs",
    primary="#1B365D",
    secondary="#A7C4E0",
    background="#EDEDED",
    footer_text="Investment Banking",
    bank_label="Advisory\nGroup",
)

MOELIS_STYLE = BankStyle(
    name="Moelis",
    primary="#1B2A4A",
    secondary="#8B9DC3",
    background="#FFFFFF",
    confidential_text="STRICTLY CONFIDENTIAL",
    footer_text="Confidential",
    bank_label="Advisory Group",
)


# ═══════════════════════════════════════════════
# LAYOUT ENGINE — Computes spatial placement
# ═══════════════════════════════════════════════

@dataclass
class Element:
    """A positioned element ready for rendering."""
    type: str  # "text", "shape", "table", "badge"
    left: float   # inches
    top: float
    width: float
    height: float
    props: dict = field(default_factory=dict)


class LayoutEngine:
    """Computes spatial layout from content specs using design rules."""

    def __init__(self, style: BankStyle):
        self.style = style
        self.slide_w = 13.333
        self.slide_h = 7.5

    @property
    def content_left(self): return self.style.margin_left
    @property
    def content_top(self): return self.style.margin_top
    @property
    def content_width(self): return self.slide_w - self.style.margin_left - self.style.margin_right
    @property
    def content_bottom(self): return self.slide_h - self.style.margin_bottom
    @property
    def content_height(self): return self.content_bottom - self.content_top

    def layout_cover(self, spec: dict) -> list[Element]:
        """Cover slide layout."""
        elements = []
        s = self.style

        # Logo placeholder
        elements.append(Element("logo", 0.45, 0.25, 1.0, 0.6, {}))

        # Confidential
        elements.append(Element("confidential", 11.5, 0.2, 1.5, 0.3, {}))

        # Title — vertically centered, left-aligned
        elements.append(Element("text", 0.5, 2.6, 9.0, 1.4, {
            "text": spec.get("title", "Discussion Materials"),
            "size": 36, "bold": True, "color": s.primary
        }))

        # Subtitle
        elements.append(Element("text", 0.5, 4.2, 5.0, 0.4, {
            "text": spec.get("subtitle", s.bank_label + " LLC"),
            "size": 18, "color": s.primary, "alpha": 0.6
        }))

        # Date
        elements.append(Element("text", 0.5, 6.5, 3.0, 0.3, {
            "text": spec.get("date", ""), "size": 12, "color": s.text_muted
        }))

        # Badge
        if s.has_badge:
            elements.append(Element("badge", 10.3, 6.85, s.badge_width, 0.35, {}))

        return elements

    def layout_section_divider(self, spec: dict) -> list[Element]:
        """Full navy background with centered title."""
        elements = []
        # Full background
        elements.append(Element("shape", 0, 0, self.slide_w, self.slide_h, {
            "fill": self.style.primary, "type": "rectangle"
        }))
        # Title — centered vertically, left-aligned with margin
        elements.append(Element("text", 0.8, 3.0, 8.0, 1.0, {
            "text": spec.get("title", ""),
            "size": 28, "bold": True, "color": "#FFFFFF"
        }))
        # Badge
        if self.style.has_badge:
            elements.append(Element("badge", 10.3, 6.85, self.style.badge_width, 0.35, {}))
        return elements

    def layout_financial_table(self, spec: dict) -> list[Element]:
        """Financial data table with auto-computed column widths and right-alignment."""
        elements = []
        s = self.style

        # Chrome (logo, confidential, slide number)
        elements.extend(self._chrome(spec))

        # Action title
        elements.append(Element("text", s.margin_left, 1.2, 11.0, 0.8, {
            "text": spec.get("title", "[Action Title]"),
            "size": s.title_size, "bold": True, "color": s.primary
        }))

        # TABLE LAYOUT ALGORITHM
        data = spec.get("data", {})
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header

        # Column width computation:
        # Label column = 30% of table width
        # Data columns = equal split of remaining 70%
        table_width = self.content_width * 0.82  # 82% of content area
        table_left = s.margin_left + (self.content_width - table_width) / 2  # centered
        label_width_pct = 0.30
        data_width_pct = (1.0 - label_width_pct) / max(1, num_cols - 1)

        col_widths = [table_width * label_width_pct]
        for i in range(num_cols - 1):
            col_widths.append(table_width * data_width_pct)

        # Row height: auto based on content
        header_height = 0.4
        row_height = 0.38
        table_height = header_height + row_height * len(rows)

        # Ensure table fits in content area
        table_top = 2.2
        max_table_height = self.content_bottom - table_top - 0.5  # leave room for source
        if table_height > max_table_height:
            row_height = (max_table_height - header_height) / len(rows)
            table_height = max_table_height

        elements.append(Element("table", table_left, table_top, table_width, table_height, {
            "headers": headers,
            "rows": rows,
            "col_widths": col_widths,
            "header_height": header_height,
            "row_height": row_height,
        }))

        # Source line
        elements.append(Element("text", s.margin_left, self.content_bottom - 0.1, 10.0, 0.3, {
            "text": spec.get("source", ""), "size": 8, "italic": True, "color": s.text_muted
        }))

        return elements

    def layout_sensitivity(self, spec: dict) -> list[Element]:
        """Sensitivity grid with highlighted base case."""
        elements = []
        s = self.style
        elements.extend(self._chrome(spec))

        # Action title
        elements.append(Element("text", s.margin_left, 1.2, 11.0, 0.8, {
            "text": spec.get("title", "[Action Title]"),
            "size": s.title_size, "bold": True, "color": s.primary
        }))

        # Table label
        elements.append(Element("text", s.margin_left, 2.1, 5.0, 0.3, {
            "text": spec.get("table_label", "Implied Share Price ($)"),
            "size": 12, "bold": True, "color": s.primary
        }))

        # SENSITIVITY TABLE LAYOUT
        data = spec.get("data", {})
        row_headers = data.get("row_headers", [])
        col_headers = data.get("col_headers", [])
        values = data.get("values", [])
        base_row = data.get("base_row", 0)
        base_col = data.get("base_col", 0)

        num_rows = len(row_headers) + 1
        num_cols = len(col_headers) + 1

        # Sensitivity tables should be compact
        table_width = min(self.content_width * 0.6, 7.5)
        table_left = s.margin_left
        table_top = 2.5

        elements.append(Element("sensitivity", table_left, table_top, table_width, 0, {
            "row_headers": row_headers,
            "col_headers": col_headers,
            "values": values,
            "base_row": base_row,
            "base_col": base_col,
            "corner_label": data.get("corner_label", "WACC \\ TGR"),
        }))

        # Note
        elements.append(Element("text", s.margin_left, self.content_bottom - 0.1, 10.0, 0.3, {
            "text": spec.get("note", ""), "size": 8, "italic": True, "color": s.text_muted
        }))

        return elements

    def layout_key_terms(self, spec: dict) -> list[Element]:
        """Moelis-style key terms table with pill labels."""
        elements = []
        s = self.style
        elements.extend(self._chrome(spec))

        # Title
        elements.append(Element("text", s.margin_left, 1.2, 11.0, 0.5, {
            "text": spec.get("title", "[Title]"),
            "size": s.title_size, "bold": True, "color": s.primary
        }))

        # Full-width header bar
        elements.append(Element("shape", s.margin_left, 1.9, self.content_width, 0.5, {
            "fill": s.primary, "type": "rectangle",
            "text": spec.get("header_text", ""),
            "text_color": "#FFFFFF", "text_bold": True, "text_size": 11
        }))

        # KEY TERMS LAYOUT ALGORITHM
        terms = spec.get("terms", [])
        pill_width = self.content_width * 0.18  # 18% for pill column
        content_width_t = self.content_width * 0.78  # 78% for content
        gap = self.content_width * 0.02

        y = 2.55
        row_padding = 0.15

        for i, term in enumerate(terms):
            # Compute row height based on number of bullet lines
            bullets = term.get("bullets", [])
            line_height = 0.22
            row_h = max(0.55, len(bullets) * line_height + 0.2)

            # Alternating background
            if i % 2 == 0:
                elements.append(Element("shape", s.margin_left, y - 0.05, self.content_width, row_h + 0.1, {
                    "fill": s.alt_row, "type": "rectangle"
                }))

            # Navy pill
            pill_y = y + (row_h - 0.4) / 2  # vertically center pill
            elements.append(Element("pill", s.margin_left + 0.15, pill_y, pill_width, 0.4, {
                "text": term["label"], "fill": s.primary, "text_color": "#FFFFFF"
            }))

            # Bullet content
            bullet_text = "\n".join(f"•  {b}" for b in bullets)
            elements.append(Element("text", s.margin_left + pill_width + gap + 0.15, y, content_width_t, row_h, {
                "text": bullet_text, "size": 10, "color": s.text_primary
            }))

            y += row_h + row_padding

        return elements

    def layout_exec_summary(self, spec: dict) -> list[Element]:
        """GS-style executive summary with callout bar and structured bullets."""
        elements = []
        s = self.style
        elements.extend(self._chrome(spec))

        # Title
        elements.append(Element("text", s.margin_left, 1.2, 11.0, 0.5, {
            "text": spec.get("title", "Executive Summary"),
            "size": s.title_size, "bold": True, "color": s.primary
        }))

        # Blue callout bar
        elements.append(Element("callout", s.margin_left, 2.0, self.content_width, 0.55, {
            "text": spec.get("callout", ""),
            "fill": s.header_fill, "border_color": s.primary,
            "text_size": 11, "text_bold": True, "text_italic": True
        }))

        # Bullet points
        points = spec.get("points", [])
        y = 2.8
        for point in points:
            # Main bullet (■)
            main_text = f"■  {point['main']}"
            elements.append(Element("text", s.margin_left + 0.2, y, 11.5, 0.3, {
                "text": main_text, "size": 11, "color": s.text_primary
            }))
            y += 0.35

            # Sub-bullets (—)
            for sub in point.get("subs", []):
                elements.append(Element("text", s.margin_left + 0.6, y, 11.0, 0.25, {
                    "text": f"—  {sub}", "size": 10, "color": s.text_primary
                }))
                y += 0.28
            y += 0.15

        return elements

    def _chrome(self, spec: dict) -> list[Element]:
        """Standard slide chrome: logo, confidential, slide number."""
        elements = []
        elements.append(Element("logo", 0.45, 0.25, 1.0, 0.6, {}))
        elements.append(Element("confidential", 11.5, 0.2, 1.5, 0.3, {}))
        if "slide_number" in spec:
            elements.append(Element("slide_number", 12.0, 7.0, 1.0, 0.3, {
                "number": spec["slide_number"]
            }))
        return elements


# ═══════════════════════════════════════════════
# RENDERER — Converts layout elements to PPTX
# ═══════════════════════════════════════════════

class Renderer:
    """Converts positioned elements to python-pptx objects."""

    def __init__(self, style: BankStyle):
        self.style = style
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def render_slide(self, elements: list[Element], bg_color: Optional[str] = None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        s = self.style

        # Background
        bg = bg_color or s.background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = s.rgb(bg)

        for el in elements:
            if el.type == "text":
                self._render_text(slide, el)
            elif el.type == "shape":
                self._render_shape(slide, el)
            elif el.type == "table":
                self._render_financial_table(slide, el)
            elif el.type == "sensitivity":
                self._render_sensitivity(slide, el)
            elif el.type == "badge":
                self._render_badge(slide, el)
            elif el.type == "logo":
                self._render_logo(slide, el)
            elif el.type == "confidential":
                self._render_confidential(slide, el)
            elif el.type == "slide_number":
                self._render_slide_number(slide, el)
            elif el.type == "pill":
                self._render_pill(slide, el)
            elif el.type == "callout":
                self._render_callout(slide, el)

    def _render_text(self, slide, el: Element):
        box = slide.shapes.add_textbox(Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = el.props.get("text", "")
        p.font.size = Pt(el.props.get("size", 11))
        p.font.bold = el.props.get("bold", False)
        p.font.italic = el.props.get("italic", False)
        p.font.color.rgb = self.style.rgb(el.props.get("color", self.style.text_primary))
        p.font.name = self.style.font
        p.alignment = el.props.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(0)
        p.space_before = Pt(0)

    def _render_shape(self, slide, el: Element):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.style.rgb(el.props.get("fill", self.style.primary))
        shape.line.fill.background()

        if "text" in el.props and el.props["text"]:
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = el.props["text"]
            p.font.size = Pt(el.props.get("text_size", 11))
            p.font.bold = el.props.get("text_bold", False)
            p.font.color.rgb = self.style.rgb(el.props.get("text_color", "#FFFFFF"))
            p.alignment = PP_ALIGN.CENTER

    def _render_financial_table(self, slide, el: Element):
        s = self.style
        headers = el.props["headers"]
        rows = el.props["rows"]
        num_cols = len(headers)
        num_rows = len(rows) + 1

        tbl_shape = slide.shapes.add_table(num_rows, num_cols,
            Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        tbl = tbl_shape.table

        # Set column widths
        col_widths = el.props.get("col_widths", [el.width / num_cols] * num_cols)
        for c, w in enumerate(col_widths):
            tbl.columns[c].width = Inches(w)

        # Header row
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(s.table_size)
            p.font.bold = True
            p.font.color.rgb = s.rgb("#FFFFFF")
            p.font.name = s.font
            p.alignment = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = s.rgb(s.primary)
            # Vertical centering
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for r, row_data in enumerate(rows):
            style_type = row_data.get("style", "normal")
            for c, val in enumerate(row_data.get("values", [])):
                cell = tbl.cell(r + 1, c)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(s.table_size)
                p.font.name = s.font
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # RIGHT-ALIGN all numeric columns (c > 0) — THE KEY RULE
                p.alignment = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT

                if style_type == "bold":
                    p.font.bold = True
                    p.font.color.rgb = s.rgb(s.text_primary)
                elif style_type == "italic_gray":
                    p.font.italic = True
                    p.font.color.rgb = s.rgb(s.text_muted)
                    if c == 0:
                        cell.text = f"  {val}"  # indent
                else:
                    p.font.color.rgb = s.rgb(s.text_primary)

                # Alternating row fill
                if r % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = s.rgb(s.alt_row)

    def _render_sensitivity(self, slide, el: Element):
        s = self.style
        p = el.props
        row_hdrs = p["row_headers"]
        col_hdrs = p["col_headers"]
        values = p["values"]
        base_r = p.get("base_row", -1)
        base_c = p.get("base_col", -1)

        num_rows = len(row_hdrs) + 1
        num_cols = len(col_hdrs) + 1

        # Auto-compute table dimensions
        cell_w = el.width / num_cols
        cell_h = 0.38
        table_h = cell_h * num_rows

        tbl_shape = slide.shapes.add_table(num_rows, num_cols,
            Inches(el.left), Inches(el.top), Inches(el.width), Inches(table_h))
        tbl = tbl_shape.table

        # Corner cell
        cell = tbl.cell(0, 0)
        cell.text = p.get("corner_label", "")
        cp = cell.text_frame.paragraphs[0]
        cp.font.size = Pt(9); cp.font.bold = True; cp.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.header_fill)

        # Column headers
        for c, hdr in enumerate(col_hdrs):
            cell = tbl.cell(0, c + 1)
            cell.text = hdr
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(9); cp.font.bold = True; cp.alignment = PP_ALIGN.CENTER
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.header_fill)

        # Row headers + data
        for r, (rhdr, row_vals) in enumerate(zip(row_hdrs, values)):
            cell = tbl.cell(r + 1, 0)
            cell.text = rhdr
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(9); cp.font.bold = True; cp.alignment = PP_ALIGN.CENTER
            cell.fill.solid(); cell.fill.fore_color.rgb = s.rgb(s.header_fill)

            for c, val in enumerate(row_vals):
                cell = tbl.cell(r + 1, c + 1)
                cell.text = str(val)
                cp = cell.text_frame.paragraphs[0]
                cp.font.size = Pt(9); cp.alignment = PP_ALIGN.CENTER
                # Highlight base case
                if r == base_r and c == base_c:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = s.rgb(s.highlight)
                    cp.font.bold = True

    def _render_badge(self, slide, el: Element):
        s = self.style
        # Navy rectangle
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(el.left), Inches(el.top), Inches(s.badge_width), Inches(el.height))
        rect.fill.solid(); rect.fill.fore_color.rgb = s.rgb(s.primary)
        rect.line.fill.background()
        tf = rect.text_frame
        p = tf.paragraphs[0]
        p.text = s.footer_text; p.font.size = Pt(10)
        p.font.bold = True; p.font.color.rgb = s.rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER
        # Light strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(el.left + s.badge_width), Inches(el.top),
            Inches(s.badge_strip_width), Inches(el.height))
        strip.fill.solid(); strip.fill.fore_color.rgb = s.rgb(s.secondary)
        strip.line.fill.background()

    def _render_logo(self, slide, el: Element):
        s = self.style
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        box.fill.background()
        box.line.color.rgb = s.rgb(s.primary)
        box.line.width = Pt(1)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = s.bank_label; p.font.size = Pt(10)
        p.font.bold = True; p.font.color.rgb = s.rgb(s.primary)
        p.alignment = PP_ALIGN.CENTER

    def _render_confidential(self, slide, el: Element):
        s = self.style
        box = slide.shapes.add_textbox(Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        p = box.text_frame.paragraphs[0]
        p.text = s.confidential_text; p.font.size = Pt(11)
        p.font.italic = True; p.font.color.rgb = s.rgb(s.confidential_color)
        p.alignment = PP_ALIGN.RIGHT

    def _render_slide_number(self, slide, el: Element):
        s = self.style
        box = slide.shapes.add_textbox(Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        p = box.text_frame.paragraphs[0]
        p.text = f"{s.footer_text} | {el.props['number']}"
        p.font.size = Pt(s.footer_size); p.font.color.rgb = s.rgb(s.text_muted)
        p.alignment = PP_ALIGN.RIGHT

    def _render_pill(self, slide, el: Element):
        s = self.style
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = s.rgb(el.props.get("fill", s.primary))
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = el.props.get("text", "")
        p.font.size = Pt(9); p.font.bold = True
        p.font.color.rgb = s.rgb(el.props.get("text_color", "#FFFFFF"))
        p.alignment = PP_ALIGN.CENTER

    def _render_callout(self, slide, el: Element):
        s = self.style
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(el.left), Inches(el.top), Inches(el.width), Inches(el.height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = s.rgb(el.props.get("fill", s.header_fill))
        shape.line.color.rgb = s.rgb(el.props.get("border_color", s.primary))
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = el.props.get("text", "")
        p.font.size = Pt(el.props.get("text_size", 11))
        p.font.bold = el.props.get("text_bold", False)
        p.font.italic = el.props.get("text_italic", False)
        p.font.color.rgb = s.rgb(s.primary)
        p.alignment = PP_ALIGN.LEFT

    def save(self, path: str):
        self.prs.save(path)
