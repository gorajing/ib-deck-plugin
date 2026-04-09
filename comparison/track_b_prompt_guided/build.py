#!/usr/bin/env python3
"""
Track B: Prompt-guided spatial code.

Represents what an LLM would write when asked to build a financial summary slide
directly with python-pptx, following the prompt in ../input/prompt.txt and IB
formatting conventions.

This is a single good-faith implementation — not engineered to fail, not
engineered to succeed. It's what I would write on a first pass given the prompt
and no reference to an existing renderer. Every Inches(), Pt(), and alignment
decision is made explicitly in this file because the prompt-guided architecture
requires the LLM to make them.

The point of the comparison is not to say this file is bad. It's to show how
much spatial code is required to produce the same slide, and to make the
tradeoffs between the two architectures visible.

Run:
    python3 build.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))

# ---------- Fixed data from the prompt ----------

TITLE = "Addus Has Delivered Consistent Revenue Growth with Expanding Profitability"
SECTION_HEADER = "Historical Financial Summary ($ in thousands)"
HEADERS = ["Metric", "FY2023A", "FY2024A", "FY2025A", "3Y CAGR"]

ROWS = [
    {"label": "Revenue",      "values": ["1,058,651", "1,154,599", "1,422,530", "14.4%"], "style": "bold"},
    {"label": "  % Growth",   "values": ["11.3%",     "9.1%",      "23.2%",     ""],      "style": "pct"},
    {"label": "Gross Profit", "values": ["339,876",   "375,021",   "461,874",   "16.6%"], "style": "bold"},
    {"label": "  % Margin",   "values": ["32.1%",     "32.5%",     "32.5%",     ""],      "style": "pct"},
    {"label": "EBITDA",       "values": ["105,082",   "116,221",   "155,027",   "21.4%"], "style": "highlight"},
    {"label": "  % Margin",   "values": ["9.9%",      "10.1%",     "10.9%",     ""],      "style": "pct"},
    {"label": "Net Income",   "values": ["62,516",    "73,598",    "95,910",    "23.8%"], "style": "bold"},
    {"label": "Diluted EPS",  "values": ["$3.83",     "$4.23",     "$5.22",     "16.7%"], "style": "normal"},
]

SOURCE = "Source: SEC EDGAR. ADUS FY2025 10-K filed February 24, 2026."
SLIDE_NUMBER = 6

# ---------- Colors ----------

NAVY = RGBColor(0x1B, 0x36, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG = RGBColor(0xED, 0xED, 0xED)
MUTED = RGBColor(0x66, 0x66, 0x66)
BLACK = RGBColor(0x1F, 0x1F, 0x1F)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
RED = RGBColor(0xC0, 0x00, 0x00)

# ---------- Build the slide ----------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = GRAY_BG

# "Confidential" top-right
box = slide.shapes.add_textbox(Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.3))
p = box.text_frame.paragraphs[0]
p.text = "Confidential"
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = RED
p.alignment = PP_ALIGN.RIGHT

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.8))
p = title_box.text_frame.paragraphs[0]
p.text = TITLE
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = NAVY

# Navy section header bar
header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(2.1),
    Inches(12.333), Inches(0.35),
)
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = NAVY
header_bar.line.fill.background()
p = header_bar.text_frame.paragraphs[0]
p.text = SECTION_HEADER
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE

# Table
n_rows = len(ROWS) + 1  # +1 for header
n_cols = len(HEADERS)
table_top = 2.55
table_left = 0.5
table_width = 12.333
row_height = 0.32

tbl_shape = slide.shapes.add_table(
    n_rows, n_cols,
    Inches(table_left), Inches(table_top),
    Inches(table_width), Inches(n_rows * row_height),
)
tbl = tbl_shape.table

# Column widths: label column wide, data columns equal
label_col_w = 2.8
data_col_w = (table_width - label_col_w) / (n_cols - 1)
tbl.columns[0].width = Inches(label_col_w)
for c in range(1, n_cols):
    tbl.columns[c].width = Inches(data_col_w)

# Header row
for c, hdr in enumerate(HEADERS):
    cell = tbl.cell(0, c)
    cell.text = hdr
    cp = cell.text_frame.paragraphs[0]
    cp.font.size = Pt(10)
    cp.font.bold = True
    cp.font.color.rgb = NAVY
    cp.alignment = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT

# Data rows
for r, row in enumerate(ROWS):
    style = row["style"]

    # Label cell
    cell = tbl.cell(r + 1, 0)
    cell.text = row["label"]
    cp = cell.text_frame.paragraphs[0]
    cp.font.size = Pt(10)
    if style == "bold":
        cp.font.bold = True
    elif style == "pct":
        cp.font.italic = True
        cp.font.color.rgb = MUTED
    elif style == "highlight":
        cp.font.bold = True
        cp.font.color.rgb = NAVY
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE

    # Data value cells
    for c, val in enumerate(row["values"]):
        cell = tbl.cell(r + 1, c + 1)
        cell.text = val
        cp = cell.text_frame.paragraphs[0]
        cp.font.size = Pt(10)
        cp.alignment = PP_ALIGN.RIGHT
        if style == "bold":
            cp.font.bold = True
        elif style == "pct":
            cp.font.italic = True
            cp.font.color.rgb = MUTED
        elif style == "highlight":
            cp.font.bold = True
            cp.font.color.rgb = NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE

# Source at the bottom
source_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(6.5),
    Inches(12), Inches(0.3),
)
p = source_box.text_frame.paragraphs[0]
p.text = SOURCE
p.font.size = Pt(8)
p.font.italic = True
p.font.color.rgb = MUTED

# Slide number in footer
num_box = slide.shapes.add_textbox(
    Inches(12.5), Inches(7.05),
    Inches(0.7), Inches(0.3),
)
p = num_box.text_frame.paragraphs[0]
p.text = str(SLIDE_NUMBER)
p.font.size = Pt(9)
p.font.color.rgb = MUTED
p.alignment = PP_ALIGN.RIGHT

# Save three runs
for run_num in range(1, 4):
    output = os.path.join(HERE, f"run{run_num}.pptx")
    prs.save(output)
    print(f"✓ Wrote {output}")

# Count spatial code decisions
print("\nTrack B complete. 3 runs produced.")
print("Lines in this file:", sum(1 for _ in open(__file__)))
print("Explicit Inches() / Pt() / RGBColor() calls below:")
with open(__file__) as f:
    content = f.read()
print(f"  Inches(): {content.count('Inches(')}")
print(f"  Pt(): {content.count('Pt(')}")
print(f"  RGBColor(): {content.count('RGBColor(')}")
print(f"  PP_ALIGN: {content.count('PP_ALIGN.')}")
